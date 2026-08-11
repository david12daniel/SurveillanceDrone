# -*- coding: utf-8 -*-
"""
Critical Design Review (CDR) deck generator — Thermal Surveillance UAS.
Structured to the standard design-review content areas (MIL-STD-1521B App. E /
IEEE 15288.2). Embeds the SysML diagram images, the requirements-traceability
CSV, and the four engineering analyses (endurance, thermal, RF, cost/fit).
"""
import csv, os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from PIL import Image

# ----------------------------------------------------------------------------
# Paths
# ----------------------------------------------------------------------------
PROJ = r"c:\Users\Josiah Laperriere\Documents\Coding\SurveillanceDrone\SurveillanceDrone"
IMG  = os.path.join(PROJ, "presentation", "assets", "diagrams")
COMP = os.path.join(PROJ, "presentation", "assets", "components")
CSV  = os.path.join(PROJ, "analysis", "requirements_traceability.csv")
COSTCHART = os.path.join(PROJ, "analysis", "cost_vs_flighttime.png")
OUT  = os.path.join(PROJ, "CDR_Thermal_Surveillance_Drone.pptx")

def im(name): return os.path.join(IMG, name)
def comp(name): return os.path.join(COMP, name)

# ----------------------------------------------------------------------------
# Theme
# ----------------------------------------------------------------------------
NAVY   = RGBColor(0x0F, 0x25, 0x40)
SLATE  = RGBColor(0x1F, 0x3A, 0x5F)
ACCENT = RGBColor(0x2E, 0x86, 0xAB)
AMBER  = RGBColor(0xD9, 0x82, 0x1E)
LIGHT  = RGBColor(0xEE, 0xF2, 0xF6)
PANEL  = RGBColor(0xF6, 0xF8, 0xFA)
TEXT   = RGBColor(0x1E, 0x28, 0x33)
MUTED  = RGBColor(0x5C, 0x6B, 0x7A)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
OKGRN  = RGBColor(0x2E, 0x7D, 0x32)
RED    = RGBColor(0xC1, 0x27, 0x27)
LINEC  = RGBColor(0xD4, 0xDD, 0xE4)

FONT   = "Segoe UI"
FONTL  = "Segoe UI Light"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

ML = Inches(0.62)          # left margin
CW = Inches(12.09)         # content width
slide_no = 0

# ----------------------------------------------------------------------------
# Low-level helpers
# ----------------------------------------------------------------------------
def _set_fill(shape, color):
    shape.fill.solid(); shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def rect(slide, l, t, w, h, color, line=None, line_w=None, shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, l, t, w, h)
    sp.shadow.inherit = False
    if color is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = line_w or Pt(1)
    return sp

def txt(slide, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        space_after=Pt(4), line_spacing=1.0, wrap=True):
    """runs: list of paragraphs; each paragraph is a list of (text, size, bold, color, italic) tuples
       OR a dict {'runs':[...], 'align':, 'level':, 'bullet':, 'space_before':, 'space_after':}."""
    tb = slide.shapes.add_textbox(l, t, w, h); tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    first = True
    for para in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        cfg = {}
        if isinstance(para, dict):
            cfg = para; rlist = para['runs']
        else:
            rlist = para
        p.alignment = cfg.get('align', align)
        p.space_after = cfg.get('space_after', space_after)
        p.space_before = cfg.get('space_before', Pt(0))
        p.line_spacing = cfg.get('line_spacing', line_spacing)
        if 'level' in cfg: p.level = cfg['level']
        for r in rlist:
            text, size, bold, color = r[0], r[1], r[2], r[3]
            italic = r[4] if len(r) > 4 else False
            run = p.add_run(); run.text = text
            run.font.name = FONT; run.font.size = Pt(size); run.font.bold = bold
            run.font.italic = italic; run.font.color.rgb = color
    return tb

def footer(slide, title_txt="Thermal Surveillance UAS  ·  Critical Design Review"):
    global slide_no
    slide_no += 1
    rect(slide, 0, SH - Inches(0.34), SW, Pt(1.2), LINEC)
    txt(slide, ML, SH - Inches(0.32), Inches(9), Inches(0.28),
        [[(title_txt, 8.5, False, MUTED)]], anchor=MSO_ANCHOR.MIDDLE)
    txt(slide, SW - Inches(2.3), SH - Inches(0.32), Inches(1.68), Inches(0.28),
        [[("CDR  ·  2026-07-26", 8.5, False, MUTED)]], align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    txt(slide, SW - Inches(0.62), SH - Inches(0.32), Inches(0.4), Inches(0.28),
        [[(str(slide_no), 8.5, True, SLATE)]], align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

def header(slide, kicker, title):
    """Standard content-slide header: eyebrow kicker + title + accent rule."""
    txt(slide, ML, Inches(0.34), CW, Inches(0.26),
        [[(kicker.upper(), 11, True, ACCENT)]])
    txt(slide, ML, Inches(0.58), CW, Inches(0.62),
        [[(title, 25, True, NAVY)]])
    rect(slide, ML, Inches(1.20), Inches(0.9), Pt(3), AMBER)

def content_slide(kicker, title):
    s = prs.slides.add_slide(BLANK)
    header(s, kicker, title)
    footer(s)
    return s

def fit_box(img_path, box_l, box_t, box_w, box_h):
    """Return (l,t,w,h) to fit image in box, centered, preserving aspect."""
    iw, ih = Image.open(img_path).size
    ar = iw / ih
    bw, bh = box_w, box_h
    if bw / bh > ar:      # box wider -> height-bound
        h = bh; w = int(bh * ar)
    else:                 # box taller -> width-bound
        w = bw; h = int(bw / ar)
    l = box_l + (bw - w) // 2
    t = box_t + (bh - h) // 2
    return int(l), int(t), int(w), int(h)

def framed_image(slide, img_path, box_l, box_t, box_w, box_h, pad=Inches(0.06)):
    l, t, w, h = fit_box(img_path, box_l + pad, box_t + pad, box_w - 2*pad, box_h - 2*pad)
    # white card + border sized to fitted image
    card = rect(slide, l - pad, t - pad, w + 2*pad, h + 2*pad, WHITE, line=LINEC, line_w=Pt(1))
    slide.shapes.add_picture(img_path, l, t, w, h)
    return l, t, w, h

def caption(slide, l, t, w, text):
    txt(slide, l, t, w, Inches(0.5),
        [[(text, 10.5, False, MUTED, True)]], align=PP_ALIGN.CENTER)

def set_cell(cell, text, size=10, bold=False, color=TEXT, fill=None,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE):
    cell.margin_left = Inches(0.06); cell.margin_right = Inches(0.06)
    cell.margin_top = Inches(0.02); cell.margin_bottom = Inches(0.02)
    cell.vertical_anchor = anchor
    if fill is not None:
        cell.fill.solid(); cell.fill.fore_color.rgb = fill
    else:
        cell.fill.solid(); cell.fill.fore_color.rgb = WHITE
    tf = cell.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    # allow embedded newlines
    parts = str(text).split("\n")
    for i, seg in enumerate(parts):
        pp = p if i == 0 else tf.add_paragraph()
        pp.alignment = align
        r = pp.add_run(); r.text = seg
        r.font.name = FONT; r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = color

def add_table(slide, l, t, w, rows, col_w, data, header_fill=NAVY, header_color=WHITE,
              header_size=10.5, body_size=9.5, row_h=Inches(0.3), aligns=None,
              zebra=True, body_colors=None):
    nrows = len(data); ncols = len(data[0])
    gtable = slide.shapes.add_table(nrows, ncols, l, t, w, row_h * nrows)
    tbl = gtable.table
    tbl.first_row = False; tbl.horz_banding = False
    # remove default style banding by using our own fills
    for ci, cwd in enumerate(col_w):
        tbl.columns[ci].width = cwd
    for ri in range(nrows):
        tbl.rows[ri].height = row_h if ri else Inches(0.34)
        for ci in range(ncols):
            cell = tbl.cell(ri, ci)
            al = PP_ALIGN.LEFT if aligns is None else aligns[ci]
            if ri == 0:
                set_cell(cell, data[ri][ci], size=header_size, bold=True,
                         color=header_color, fill=header_fill, align=al)
            else:
                fill = WHITE
                if zebra and ri % 2 == 0:
                    fill = PANEL
                bc = TEXT
                if body_colors and body_colors.get((ri, ci)):
                    bc = body_colors[(ri, ci)]
                set_cell(cell, data[ri][ci], size=body_size, color=bc, fill=fill, align=al)
    return tbl

def chip(slide, l, t, w, h, text, fill, tcolor=WHITE, size=10):
    c = rect(slide, l, t, w, h, fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    try: c.adjustments[0] = 0.5
    except Exception: pass
    tf = c.text_frame; tf.word_wrap = True
    tf.margin_left=Inches(0.05); tf.margin_right=Inches(0.05)
    tf.margin_top=0; tf.margin_bottom=0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.name=FONT; r.font.size=Pt(size); r.font.bold=True; r.font.color.rgb=tcolor
    return c

def stat(slide, l, t, w, h, big, label, accent=ACCENT):
    rect(slide, l, t, w, h, PANEL, line=LINEC, line_w=Pt(0.75), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(slide, l, t, Inches(0.09), h, accent)
    txt(slide, l+Inches(0.2), t+Inches(0.12), w-Inches(0.3), Inches(0.5),
        [[(big, 22, True, NAVY)]], anchor=MSO_ANCHOR.TOP)
    txt(slide, l+Inches(0.2), t+Inches(0.6), w-Inches(0.3), h-Inches(0.62),
        [[(label, 10, False, MUTED)]], anchor=MSO_ANCHOR.TOP)

def bullet_para(text="", size=14, bold=False, color=TEXT, level=0, sb=Pt(3), sa=Pt(3),
                bullet=True, lead=None, lead_color=None):
    runs = []
    if lead:
        runs.append((lead, size, True, lead_color or NAVY))
        runs.append((text, size, bold, color))
    else:
        runs.append((text, size, bold, color))
    d = {'runs': runs, 'level': level, 'space_before': sb, 'space_after': sa,
         'line_spacing': 1.04}
    return d

def bullets(slide, l, t, w, h, items, size=14):
    """items: list of dicts from bullet_para; renders a manual bullet glyph."""
    tb = slide.shapes.add_textbox(l, t, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    first = True
    for it in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
        p.space_before = it.get('space_before', Pt(3))
        p.space_after = it.get('space_after', Pt(3))
        p.line_spacing = it.get('line_spacing', 1.04)
        lvl = it.get('level', 0)
        glyph = "" if not it.get('bullet', True) else ("▸  " if lvl == 0 else "–  ")
        indent = "     " * lvl
        if glyph:
            r = p.add_run(); r.text = indent + glyph
            r.font.name=FONT; r.font.size=Pt(size); r.font.bold=True
            r.font.color.rgb = ACCENT if lvl == 0 else MUTED
        for (text, sz, bold, color) in [ (x[0],x[1],x[2],x[3]) for x in it['runs'] ]:
            r = p.add_run(); r.text = text
            r.font.name=FONT; r.font.size=Pt(sz); r.font.bold=bold; r.font.color.rgb=color
    return tb

def r(text, size, bold, color): return (text, size, bold, color)

# ============================================================================
# CSV LOAD
# ============================================================================
rows = []
with open(CSV, newline='', encoding='utf-8') as f:
    for row in csv.DictReader(f):
        rows.append(row)

MISSION_IDS = ["R1","R2","R3","R3_1","R3_2","R4","R5","R6","R7","R8"]
mission = {row['requirement_id']: row for row in rows if row['requirement_id'] in MISSION_IDS}
derived = [row for row in rows if row['requirement_id'] not in MISSION_IDS]

# ============================================================================
# 1. TITLE
# ============================================================================
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, 0, SW, Inches(0.14), AMBER)
# subtle accent block
rect(s, 0, Inches(4.9), SW, Inches(0.02), RGBColor(0x2a,0x44,0x63))
txt(s, ML, Inches(1.5), CW, Inches(0.4),
    [[("CRITICAL DESIGN REVIEW  (CDR)", 15, True, RGBColor(0x8F,0xC5,0xE0))]])
txt(s, ML, Inches(2.0), Inches(11.6), Inches(1.8),
    [[("Thermal Surveillance UAS", 46, True, WHITE)],
     [("Autonomous Wildlife-Scouting Drone", 26, False, RGBColor(0xC9,0xD8,0xE6))]],
    line_spacing=1.05)
txt(s, ML, Inches(5.15), CW, Inches(1.4),
    [[("Detailed design baseline, requirements traceability, and verifying analyses", 15, False, RGBColor(0xB9,0xCB,0xDb))],
     [("Model-Based Systems Engineering (SysML v2)  ·  Content per MIL-STD-1521B App. E / IEEE 15288.2 CDR", 12.5, False, RGBColor(0x8A,0x9E,0xB2))]],
    space_after=Pt(6))
# meta line
txt(s, ML, Inches(6.55), CW, Inches(0.4),
    [[("Review date  2026-07-26        Design maturity  Detailed design complete (pre-fabrication)        Budget  ≤ $2,500 (R4)", 11, False, RGBColor(0x7F,0x95,0xAB))]])

# ============================================================================
# 2. AGENDA / PURPOSE
# ============================================================================
s = content_slide("Section 0 · Introduction", "CDR Purpose & Agenda")
bullets(s, ML, Inches(1.45), Inches(6.7), Inches(5.2), [
    bullet_para(lead="Purpose.  ", size=14, text="Demonstrate that the detailed design meets the mission and derived requirements, that interfaces are defined, and that the design is verified by analysis and ready to proceed to fabrication / integration."),
    bullet_para(lead="Scope.  ", size=14, text="The committed Phase 1–3 system: airframe & propulsion, power, thermal payload + onboard compute, RF links, ground control, and the autonomy behavior. Phase 4 (video downlink) is shown as a deferred capability."),
    bullet_para(lead="Authoritative source.  ", size=14, text="A single SysML v2 model (requirements · architecture · behavior · analysis) under git version control, with real, purchasable components traded against the requirements."),
    bullet_para("Exit criteria (assessed on the closing slide): design traces to requirements, analyses show compliance with margin, interfaces closed, risks identified with mitigations, open items bounded.", size=13.5, color=MUTED),
])
# agenda card
rect(s, Inches(7.7), Inches(1.45), Inches(5.0), Inches(5.15), PANEL, line=LINEC, line_w=Pt(1), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, Inches(7.95), Inches(1.62), Inches(4.6), Inches(0.4), [[("AGENDA", 12, True, ACCENT)]])
agenda = [
    "1   Requirements & traceability",
    "2   System & subsystem architecture",
    "3   Interface design",
    "4   Functional / behavioral design",
    "5   Engineering analyses (endurance,\n      thermal, RF, fit, cost)",
    "6   Verification & validation plan",
    "7   Manufacturing & producibility",
    "8   Risk, open items & CDR closure",
]
txt(s, Inches(7.95), Inches(2.1), Inches(4.55), Inches(4.4),
    [[(a, 13.5, False, TEXT)] for a in agenda], space_after=Pt(9), line_spacing=1.0)

# ============================================================================
# 3. CDR CONTENT COMPLIANCE (DID mapping)
# ============================================================================
s = content_slide("Section 0 · Introduction", "CDR Data-Item Content — Coverage Map")
txt(s, ML, Inches(1.28), CW, Inches(0.4),
    [[("Standard design-review content areas and where each is addressed in this package.", 12.5, False, MUTED)]])
data = [["CDR content area (per DID / MIL-STD-1521B App. E)", "Addressed in", "Status"]]
cov = [
    ("Mission / operational concept (ConOps)", "Mission & ConOps", "Complete"),
    ("System & performance requirements + traceability", "Requirements (RTM)", "Complete"),
    ("Detailed design — physical & functional architecture", "Architecture · Behavior", "Complete"),
    ("Interface design (internal & external / IRS)", "Interface Design", "Complete"),
    ("Design & specialty-engineering analyses", "Analyses (×4)", "Complete"),
    ("Trade studies / component selection", "Design Selection", "Complete"),
    ("Verification & validation approach (VCRM)", "V&V Plan", "Plan defined"),
    ("Producibility / manufacturing", "Manufacturing", "Complete"),
    ("Reliability & safety (failsafes)", "State machine · Risk", "Complete"),
    ("Risk assessment & open items (TBD/TBR)", "Risk · Open Items", "Complete"),
]
for a,b,c in cov: data.append([a,b,c])
bc = {}
for i in range(1, len(data)):
    st = data[i][2]
    bc[(i,2)] = OKGRN if st=="Complete" else AMBER
add_table(s, ML, Inches(1.7), CW, len(data),
          [Inches(7.09), Inches(3.0), Inches(2.0)], data,
          aligns=[PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.CENTER],
          row_h=Inches(0.42), body_size=11, header_size=11, body_colors=bc)

# ============================================================================
# 4. MISSION & CONOPS
# ============================================================================
s = content_slide("Section 1 · System Overview", "Mission & Concept of Operations")
bullets(s, ML, Inches(1.4), Inches(6.7), Inches(4.2), [
    bullet_para(lead="Mission.  ", size=14, text="Autonomous daytime aerial survey to detect and classify wildlife (deer, turkey, other animals) and humans by their thermal (IR) signature over open / semi-open terrain."),
    bullet_para(lead="Operating point.  ", size=14, text="90–120 m AGL (R1), 2.23 m/s ground speed (R2), ≥ 2.8 km survey line (R7) in ≤ 4.5 m/s sustained wind, with ≥ 5 °C target-to-background thermal contrast."),
    bullet_para(lead="Sortie.  ", size=14, text="Plan waypoint route → auto takeoff → cruise-survey with live onboard thermal inference → investigate & classify detections → return-to-launch & land. Failsafes on link-loss and low battery."),
    bullet_para(lead="Autonomy stance.  ", size=14, text="Thermal streams live to the onboard computer for real-time inference that drives route changes — no recording and no required downlink in the committed build; a thermal-video downlink is a planned future capability (Phase 4)."),
])
# ConOps flow image (tall) on right
framed_image(s, im("conduct_sortie_action_flow.png"), Inches(7.7), Inches(1.35), Inches(2.4), Inches(5.25))
caption(s, Inches(7.6), Inches(6.62), Inches(2.6), "ConductSortie — top-level action flow")
# key params stats column
stat(s, Inches(10.35), Inches(1.35), Inches(2.35), Inches(1.22), "90–120 m", "Survey altitude AGL (R1)")
stat(s, Inches(10.35), Inches(2.70), Inches(2.35), Inches(1.22), "2.8 km", "Survey range in wind (R7)", accent=AMBER)
stat(s, Inches(10.35), Inches(4.05), Inches(2.35), Inches(1.22), "≥ 30 min", "Endurance, 60 stretch (R6/R8)")
stat(s, Inches(10.35), Inches(5.40), Inches(2.35), Inches(1.22), "≤ $2,500", "Total system cost (R4)", accent=OKGRN)

# ============================================================================
# Section divider helper
# ============================================================================
def section_divider(num, title, subtitle=""):
    global slide_no
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, SW, SH, NAVY)
    rect(s, 0, 0, Inches(0.16), SH, AMBER)
    rect(s, Inches(0.62), Inches(3.02), Inches(0.55), Pt(3), ACCENT)
    txt(s, ML, Inches(2.55), CW, Inches(0.5), [[("SECTION %s" % num, 15, True, RGBColor(0x8F,0xC5,0xE0))]])
    txt(s, ML, Inches(3.15), CW, Inches(1.2), [[(title, 40, True, WHITE)]])
    if subtitle:
        txt(s, ML, Inches(4.35), Inches(10.5), Inches(1.2),
            [[(subtitle, 15, False, RGBColor(0xB9,0xCB,0xDB))]], line_spacing=1.15)
    slide_no += 1
    return s

def short(text, n=64):
    text = text.strip().strip('"')
    if len(text) <= n: return text
    cut = text[:n].rsplit(" ", 1)[0]
    return cut + "…"

# Verification method per mission requirement (I=Inspect, A=Analysis, D=Demo, T=Test)
VMETHOD = {"R1":"A, T","R2":"A, T","R3":"A, D, T","R3_1":"A, D","R3_2":"A, D, T",
           "R4":"I","R5":"I","R6":"A, T","R7":"A, T","R8":"A"}

# ============================================================================
# ---- 5. Operational Concept (OV-1) — DoDAF high-level operational graphic ----
s = content_slide("Section 1 · System Overview", "Operational Concept  (OV-1)")
framed_image(s, im("ov1_operational_concept.png"), ML, Inches(1.32), CW, Inches(5.4))
caption(s, ML, Inches(6.82), CW,
        "OV-1 High-Level Operational Concept Graphic (DoDAF) — the mission in context; callouts trace to R1 · R2 · R3 · R7")

# ============================================================================
# 6. SECTION — REQUIREMENTS
# ============================================================================
section_divider("1", "Requirements & Traceability",
    "Eight mission requirements decomposed into subsystem requirements, each allocated to a "
    "component via a formal satisfy link. Source: requirements_traceability.csv (SysML v2 model).")

# ---- Shared SysML-style diagram helpers (used by the decomposition + BDD views) ----
def vline(slide, x, y1, y2, color=ACCENT, wt=Pt(1.5)):
    rect(slide, x, min(y1,y2), wt, abs(y2-y1), color)
def hline(slide, x1, x2, y, color=ACCENT, wt=Pt(1.5)):
    rect(slide, min(x1,x2), y, abs(x2-x1), wt, color)

def sysml_box(slide, l, t, w, h, name, stereo="«part def»", fill=None, line=ACCENT,
              name_size=13, name_color=NAVY, stereo_color=MUTED):
    fill = fill if fill is not None else RGBColor(0xEA,0xF1,0xF7)
    rect(slide, l, t, w, h, fill, line=line, line_w=Pt(1.25))
    txt(slide, l, t+Inches(0.06), w, Inches(0.2), [[(stereo, 9, False, stereo_color)]], align=PP_ALIGN.CENTER)
    txt(slide, l, t+Inches(0.25), w, h-Inches(0.3), [[(name, name_size, True, name_color)]],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)

def def_panel(slide, l, t, w, h, defname, parts):
    rect(slide, l, t, w, h, WHITE, line=ACCENT, line_w=Pt(1.5))
    rect(slide, l, t, w, Inches(0.6), SLATE)
    txt(slide, l, t+Inches(0.05), w, Inches(0.2), [[("«part def»", 9, False, RGBColor(0xBB,0xD0,0xE0))]], align=PP_ALIGN.CENTER)
    txt(slide, l, t+Inches(0.24), w, Inches(0.3), [[(defname, 14, True, WHITE)]], align=PP_ALIGN.CENTER)
    avail = h - Inches(0.78)
    step = min(Inches(0.34), Emu(int(avail/max(len(parts),1))))
    yy = t + Inches(0.72)
    for role, typ in parts:
        txt(slide, l+Inches(0.2), yy, w-Inches(0.4), step,
            [[(role+"  :  ", 11.5, False, MUTED), (typ, 11.5, True, NAVY)]], anchor=MSO_ANCHOR.MIDDLE)
        yy += step

# ---- 6. Requirements decomposition tree (front of the Requirements section) ----
s = content_slide("Section 1 · Requirements", "Requirements — Decomposition by Subsystem")
txt(s, ML, Inches(1.25), CW, Inches(0.35),
    [[("Mission requirements decompose into per-subsystem requirement packages in the SysML v2 model; each subsystem requirement subsets a mission requirement and is satisfied by a component.", 12, False, MUTED)]])
sysml_box(s, Inches(4.42), Inches(1.78), Inches(4.5), Inches(0.72), "Mission Requirements  ·  R1–R8",
          stereo="«requirements»", fill=NAVY, line=NAVY, name_size=13, name_color=WHITE, stereo_color=RGBColor(0xBB,0xD0,0xE0))
subs = [("Airframe",["R4_AF"]), ("Battery",["R4_BAT"]), ("Thermal Camera",["R3_CAM"]),
        ("SBC / Compute",["R4_SBC"]), ("Ground Control",["R4_GCS"]), ("Autonomy / Behavior",["BHV"])]
def count_pref(prefs):
    return sum(1 for d in derived if any(d['requirement_id'].startswith(p) or (p=="BHV" and "_BHV_" in d['requirement_id']) for p in prefs))
n=len(subs); bw=Inches(1.83); gap=(CW-bw*n)//(n-1)
bus_y=Inches(3.2); by=Inches(3.75); bh=Inches(1.5)
vline(s, Inches(6.67), Inches(2.5), bus_y)
firstc = ML + bw//2; lastc = ML + (n-1)*(bw+gap) + bw//2
hline(s, firstc, lastc, bus_y)
for i,(nm,prefs) in enumerate(subs):
    x = ML + i*(bw+gap)
    vline(s, x+bw//2, bus_y, by)
    rect(s, x, by, bw, bh, RGBColor(0xEA,0xF1,0xF7), line=ACCENT, line_w=Pt(1.25))
    txt(s, x, by+Inches(0.12), bw, Inches(0.2), [[("«requirements»", 8.5, False, MUTED)]], align=PP_ALIGN.CENTER)
    txt(s, x, by+Inches(0.36), bw, Inches(0.55), [[(nm, 12, True, NAVY)]], align=PP_ALIGN.CENTER)
    txt(s, x, by+Inches(1.0), bw, Inches(0.35), [[("%d requirements" % count_pref(prefs), 11, True, ACCENT)]], align=PP_ALIGN.CENTER)
rect(s, ML, Inches(5.55), CW, Inches(0.95), PANEL, line=LINEC, line_w=Pt(0.75), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, ML+Inches(0.2), Inches(5.68), CW-Inches(0.4), Inches(0.75),
    [[("Purpose of this view: ", 11.5, True, NAVY),
      ("show the requirement decomposition structure — that requirements are modeled formally per subsystem and all trace up to R1–R8 and down to a component. The full requirement text and per-requirement allocation are in the three RTM slides that follow and in the model; this slide is the map, not the fine print.", 11.5, False, TEXT)]], line_spacing=1.08)

# ---- 6. Mission requirements (image + verification side table) ----
s = content_slide("Section 1 · Requirements", "Mission Requirements  (R1–R8)")
framed_image(s, im("system_level_requirements.png"), ML, Inches(1.35), Inches(7.35), Inches(5.3))
caption(s, ML, Inches(6.68), Inches(7.35), "Mission & autonomy requirements — SysML v2 Requirements Table (SysON export from model.sysml)")
# side: verification method table
vdata = [["Req", "Verify"]]
for rid in MISSION_IDS:
    vdata.append([rid.replace("_","."), VMETHOD.get(rid,"A")])
add_table(s, Inches(8.35), Inches(1.35), Inches(2.35), len(vdata),
          [Inches(1.15), Inches(1.20)], vdata,
          aligns=[PP_ALIGN.LEFT, PP_ALIGN.CENTER], row_h=Inches(0.315),
          body_size=10.5, header_size=10.5)
rect(s, Inches(10.9), Inches(1.35), Inches(1.8), Inches(4.7), PANEL, line=LINEC, line_w=Pt(0.75), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, Inches(11.05), Inches(1.5), Inches(1.55), Inches(0.4), [[("V-METHODS", 10.5, True, ACCENT)]])
txt(s, Inches(11.05), Inches(1.95), Inches(1.6), Inches(4.0),
    [[("I  Inspection", 11, False, TEXT)],
     [("A  Analysis", 11, False, TEXT)],
     [("D  Demonstration", 11, False, TEXT)],
     [("T  Test", 11, False, TEXT)],
     [("", 6, False, TEXT)],
     [("At CDR: A complete for all; T/D planned (see V&V).", 9.5, False, MUTED)]],
    space_after=Pt(8), line_spacing=1.05)

# ---- 7. Traceability approach + coverage rollup ----
s = content_slide("Section 1 · Requirements", "Traceability Approach & Coverage")
bullets(s, ML, Inches(1.35), Inches(4.7), Inches(5.2), [
    bullet_para(lead="Up-trace.  ", size=13.5, text="Each subsystem requirement subsets the mission requirement it decomposes (e.g. R3_CAM_RES ⊂ R3, cost reqs ⊂ R4)."),
    bullet_para(lead="Down-trace.  ", size=13.5, text="Each component part def carries a satisfy link to the requirements it fulfills — the RTM is derived from that web, not maintained by hand."),
    bullet_para(lead="Coverage.  ", size=13.5, text="All 8 mission requirements decompose to subsystem requirements and allocate to at least one component; behavior reqs (…_BHV_…) allocate to autonomy functions + firmware."),
    bullet_para(lead="Bidirectional.  ", size=13.5, text="Any requirement → responsible component, and any component → requirements it must meet, are both queryable from the model."),
])
# rollup table (computed) — each requirement carries a short general title for context
roll = [["Mission requirement", "Derived", "Allocated to", "Verify"]]
TITLE = {
    "R1":"Survey altitude", "R2":"Ground speed", "R3":"Thermal detect & classify",
    "R4":"System cost", "R5":"Minimize soldering", "R6":"Endurance (30 min)",
    "R7":"Range in wind", "R8":"Endurance stretch",
}
alloc_map = {
    "R1":"Airframe, autonomy (flyRoute)", "R2":"Airframe, autonomy (flyRoute)",
    "R3":"IRCamera, SBC, thermalModel", "R4":"All subsystems (cost)",
    "R5":"Airframe (PNP, pre-soldered)", "R6":"Battery, autonomy (RTL)",
    "R7":"GCS, antennas, RF links", "R8":"Battery",
}
for rid in ["R1","R2","R3","R4","R5","R6","R7","R8"]:
    n = sum(1 for d in derived if d['subsets_root'] == rid)
    n2 = sum(1 for d in derived if d['subsets_parent'] and rid in [x.strip() for x in d['subsets_parent'].split(",")])
    roll.append(["%s · %s" % (rid, TITLE[rid]), str(max(n, n2)), alloc_map[rid], VMETHOD[rid]])
add_table(s, Inches(5.55), Inches(1.35), Inches(7.16), len(roll),
          [Inches(2.95), Inches(0.9), Inches(2.4), Inches(0.91)], roll,
          aligns=[PP_ALIGN.LEFT, PP_ALIGN.CENTER, PP_ALIGN.LEFT, PP_ALIGN.CENTER],
          row_h=Inches(0.5), body_size=10, header_size=10)
txt(s, Inches(5.55), Inches(6.35), Inches(7.16), Inches(0.5),
    [[("Full matrix on the following three slides (from requirements_traceability.csv). ", 10, False, MUTED, True),
      ("Derived-req counts reflect formal subset/parent links.", 10, False, MUTED, True)]])

# ---- 8-10. RTM detail tables (split by subsystem group) ----
def rtm_slide(title, id_prefixes, note=""):
    s = content_slide("Section 1 · Requirements  ·  RTM", title)
    sub = [d for d in derived if any(d['requirement_id'].startswith(p) or (p=="BHV" and "_BHV_" in d['requirement_id']) for p in id_prefixes)]
    data = [["Requirement ID", "Requirement (abridged)", "Parent", "Allocated to"]]
    for d in sub:
        parent = d['subsets_parent'] or d['subsets_root'] or "—"
        parent = parent.replace(" ", "")
        alloc = short(d['satisfied_by'] or "—", 26)
        data.append([d['requirement_id'], short(d['doc_summary'], 70), parent, alloc])
    row_h = Inches(0.30) if len(data) <= 15 else Inches(0.275)
    add_table(s, ML, Inches(1.32), CW, len(data),
              [Inches(2.35), Inches(5.99), Inches(1.35), Inches(2.4)], data,
              aligns=[PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.CENTER, PP_ALIGN.LEFT],
              row_h=row_h, body_size=9, header_size=10)
    if note:
        txt(s, ML, SH-Inches(0.62), CW, Inches(0.3), [[(note, 9.5, False, MUTED, True)]])
    return s

rtm_slide("RTM — Airframe Subsystem & Autonomy Behavior", ["R4_AF", "BHV"],
          note="Behavior requirements (…_BHV_…) are allocated to autonomy action functions (flyRoute, executeRtl) and the flight-controller firmware (ArduPilot ArduCopter).")
rtm_slide("RTM — Battery & Camera (Thermal Payload) Subsystems", ["R4_BAT", "R3_CAM"])
rtm_slide("RTM — SBC (Onboard Compute) & Ground Control Subsystems", ["R4_SBC", "R4_GCS"])

# ---- Diagram helpers + Requirements-decomposition slide were relocated to the
#      front of the Requirements section (just after the section divider, above). ----

# ============================================================================
# SECTION 2 — ARCHITECTURE
# ============================================================================
section_divider("2", "System Architecture",
    "Physical decomposition and interfaces. The system is a drone plus a laptop-based ground "
    "control station, joined by three wireless RF links; the drone integrates propulsion, power, "
    "thermal payload, and an onboard inference computer.")

# ---- Block-definition (composition) view — SysML v2 analog of a BDD ----
s = content_slide("Section 2 · Architecture", "System Composition — Block-Definition View")
sysml_box(s, Inches(5.02), Inches(1.32), Inches(3.3), Inches(0.68), "AerialObservationSystem")
vline(s, Inches(6.67), Inches(2.0), Inches(2.2))
hline(s, Inches(3.545), Inches(9.795), Inches(2.2))
vline(s, Inches(3.545), Inches(2.2), Inches(2.34))
vline(s, Inches(9.795), Inches(2.2), Inches(2.34))
def_panel(s, Inches(0.62), Inches(2.34), Inches(5.85), Inches(4.4), "Drone", [
    ("platform","Airframe"), ("battery","Battery"), ("antiSpark","AntiSparkFilter"),
    ("camera","IRCamera"), ("fpvCam","FpvCamera"), ("gps","GpsModule"),
    ("sbc","SBCPayload"), ("rx","RadioReceiver"), ("vtx","Vtx"), ("ubec","Ubec"),
    ("vtxAntenna, rxAntennaA/B","Antenna [3]"),
])
def_panel(s, Inches(6.87), Inches(2.34), Inches(5.85), Inches(4.4), "GCS", [
    ("viewingComputer","Laptop"), ("videoRx","Vrx"), ("capture","UsbCap"),
    ("rcTx","RcTx"), ("laptopLink","TelemetryGroundLink"),
    ("groundAntenna","Antenna"), ("charger","Charger"),
])
txt(s, ML, Inches(6.82), CW, Inches(0.4),
    [[("Block-definition (composition) view — the SysML v2 analog of a BDD: each ", 10, False, MUTED),
      ("part def", 10, True, MUTED),
      (" and the typed ", 10, False, MUTED), ("part", 10, True, MUTED),
      (" usages it composes. Phase-4 OpenHD parts (air/ground Wi-Fi adapters + antennas) omitted for the committed build.", 10, False, MUTED)]])

# ---- System context IBD ----
s = content_slide("Section 2 · Architecture", "System Context — Aerial Observation System")
framed_image(s, im("full_system_internal_block_diagram.png"), ML, Inches(1.5), Inches(7.15), Inches(3.6))
caption(s, ML, Inches(5.15), Inches(7.15), "AerialObservationSystem — internal block diagram (drone ↔ GCS wireless interfaces)")
bullets(s, Inches(8.1), Inches(1.45), Inches(4.6), Inches(5.2), [
    bullet_para(lead="Two blocks + operator.  ", size=13, text="drone : Drone and gcs : GCS (laptop-based). The MacBook Air is an existing external actor, excluded from system cost."),
    bullet_para("Three wireless interfaces close the system:", size=13, bold=True, color=NAVY, bullet=True),
    bullet_para("5.8 GHz analog FPV video — drone VTX → ground patch antenna (piloting).", size=12.5, level=1),
    bullet_para("2.4 GHz ELRS — control uplink + MAVLink telemetry downlink (one link, both roles).", size=12.5, level=1),
    bullet_para("GNSS receive (satellite) for navigation.", size=12.5, level=1),
    bullet_para(lead="Wired ground interface.  ", size=13, text="Battery ↔ charger (support equipment; not flown)."),
    bullet_para(lead="Phase 4 (deferred).  ", size=13, text="Adds a 5.8 GHz OpenHD digital downlink of the thermal / AI feed — not in the committed build."),
])

# ---- 14. Drone IBD (detailed) ----
s = content_slide("Section 2 · Architecture", "Drone — Internal Block Diagram")
framed_image(s, im("drone_internal_block_diagram.png"), ML, Inches(1.3), Inches(8.4), Inches(5.35))
caption(s, ML, Inches(6.62), Inches(8.4), "Drone composition and internal wiring (power · video · data · RF)")
rect(s, Inches(9.3), Inches(1.35), Inches(3.42), Inches(5.2), PANEL, line=LINEC, line_w=Pt(0.75), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, Inches(9.5), Inches(1.5), Inches(3.1), Inches(0.35), [[("KEY ELEMENTS", 11, True, ACCENT)]])
bullets(s, Inches(9.5), Inches(1.92), Inches(3.05), Inches(4.6), [
    bullet_para(lead="platform.  ", size=11.5, text="Airframe (FC, ESC, motors, PDB) — the wiring hub.", sa=Pt(4)),
    bullet_para(lead="battery → antiSpark.  ", size=11.5, text="XT60 power chain with inrush-limiting filter.", sa=Pt(4)),
    bullet_para(lead="camera → sbc.  ", size=11.5, text="Thermal over USB-UVC to onboard compute.", sa=Pt(4)),
    bullet_para(lead="sbc ↔ FC.  ", size=11.5, text="MAVLink over UART (autonomous route mod).", sa=Pt(4)),
    bullet_para(lead="fpvCam → vtx.  ", size=11.5, text="Analog video for piloting; VTX → antenna.", sa=Pt(4)),
    bullet_para(lead="rx (diversity).  ", size=11.5, text="ELRS control + telemetry, two antennas.", sa=Pt(4)),
    bullet_para(lead="gps.  ", size=11.5, text="u-blox M10 for waypoint navigation.", sa=Pt(4)),
])

# ---- 15. SBC subsystem ----
s = content_slide("Section 2 · Architecture", "Onboard Compute (SBC) — Software Architecture")
framed_image(s, im("sbc_internal_block_diagram_dev.png"), ML, Inches(1.35), Inches(6.1), Inches(5.25))
caption(s, ML, Inches(6.62), Inches(6.1), "SBCPayload internal block diagram — amber boxes = to-be-developed software (D-1 thermalModel · D-2 missionApp)")
bullets(s, Inches(6.95), Inches(1.5), Inches(5.75), Inches(5.1), [
    bullet_para(lead="NanoPi M5 (RK3576, 6 TOPS NPU).  ", size=13, text="Passive-cooled, ≤ 10 W — meets R4_SBC_PWR / R4_SBC_TEMP with no fan required."),
    bullet_para("missionApp  (D-2 · to develop).  ", size=13, bold=True, color=AMBER, bullet=True),
    bullet_para("Receives the thermal frame (UVC), runs the detect/investigate loop, commands routes via MAVLink.", size=12.5, level=1),
    bullet_para("thermalModel  (D-1 · to develop) → rknnRuntime.  ", size=13, bold=True, color=AMBER, bullet=True),
    bullet_para("INT8 model allocated to the RKNN NPU runtime — the inference server for onboard classification.", size=12.5, level=1),
    bullet_para("mavlinkRouter.  ", size=13, bold=True, color=NAVY, bullet=True),
    bullet_para("Routes MAVLink between missionApp and the flight-controller UART (data_af interface).", size=12.5, level=1),
    bullet_para(lead="Ports.  ", size=13, text="USB-A #1 = thermal in (committed); USB-A #2 reserved for Phase 4 downlink; USB-C = 12 V power via dedicated UBEC."),
])

# ---- 16. GCS subsystem ----
s = content_slide("Section 2 · Architecture", "Ground Control Station — Internal Block Diagram")
framed_image(s, im("gcs_internal_block_diagram.png"), ML, Inches(1.35), Inches(6.1), Inches(5.25))
caption(s, ML, Inches(6.62), Inches(6.1), "GCS internal block diagram — laptop, video capture, control & telemetry")
bullets(s, Inches(6.95), Inches(1.5), Inches(5.75), Inches(5.1), [
    bullet_para(lead="viewingComputer (laptop).  ", size=13, text="Existing MacBook Air running QGroundControl 4.4+ — mission plan, telemetry, video, params, alerts, logs."),
    bullet_para(lead="videoRx → capture.  ", size=13, text="5.8 GHz analog receiver → USB UVC capture into the laptop (QuickTime/OBS) for piloting video."),
    bullet_para(lead="rcTx.  ", size=13, text="RadioMaster TX12 MkII (ELRS) — manual control uplink; primary in-flight MAVLink telemetry runs through the HGLRC Hermes USB dongle."),
    bullet_para(lead="groundAntenna.  ", size=13, text="TrueRC X-AIR patch (10 dBic, 120°) — closes the 2.8 km video link with margin."),
    bullet_para(lead="charger.  ", size=13, text="Ground-support equipment (HOTA D6 Pro) — not flown, excluded from system cost."),
    bullet_para(lead="Phase 4 (deferred) — not shown on the diagram above.  ", size=12.5, lead_color=AMBER, color=MUTED, text="Ground OpenHD node for the future thermal-video downlink: openHDRx (Alfa AWUS036ACH Wi-Fi adapter) + openHDAntennaA / openHDAntennaB (Foxeer Echo 2 Max diversity antennas). Phase 4 options, not in the committed build."),
])

# ---- 17. Interface design ----
s = content_slide("Section 2 · Architecture", "Interface Design  (Internal & External)")
txt(s, ML, Inches(1.28), Inches(6), Inches(0.3), [[("INTERNAL (WIRED) INTERFACES", 11, True, ACCENT)]])
intern = [["Interface", "From → To", "Type"]]
for a,b,c in [
    ("Battery power","battery → antiSpark → platform","XT60, 6S"),
    ("Thermal video","camera → sbc","USB-UVC"),
    ("Autonomy data","sbc ↔ FC","MAVLink / UART"),
    ("FPV video","fpvCam → vtx","Analog CVBS"),
    ("Navigation","gps → FC","UART"),
    ("Control data","rx → FC","CRSF"),
    ("SBC power","battery → UBEC → sbc","12 V / USB-C"),
]:
    intern.append([a,b,c])
add_table(s, ML, Inches(1.62), Inches(6.0), len(intern),
          [Inches(1.7), Inches(2.95), Inches(1.35)], intern,
          aligns=[PP_ALIGN.LEFT,PP_ALIGN.LEFT,PP_ALIGN.LEFT], row_h=Inches(0.34), body_size=9.5, header_size=10)

txt(s, Inches(6.9), Inches(1.28), Inches(5.8), Inches(0.3), [[("EXTERNAL (WIRELESS) INTERFACES", 11, True, ACCENT)]])
extern = [["Link", "Band", "Role"]]
for a,b,c in [
    ("Control uplink","2.4 GHz ELRS","RC + failsafe"),
    ("Telemetry downlink","2.4 GHz ELRS","MAVLink status"),
    ("FPV video","5.8 GHz analog","Piloting video"),
    ("GNSS","1.575 GHz","Position (Rx only)"),
    ("Thermal downlink*","5.8 GHz OpenHD","Phase 4 (deferred)"),
]:
    extern.append([a,b,c])
add_table(s, Inches(6.9), Inches(1.62), Inches(5.8), len(extern),
          [Inches(2.2), Inches(1.9), Inches(1.7)], extern,
          aligns=[PP_ALIGN.LEFT,PP_ALIGN.LEFT,PP_ALIGN.LEFT], row_h=Inches(0.36), body_size=9.5, header_size=10)

# compatibility rules box
rect(s, ML, Inches(4.55), CW, Inches(2.05), PANEL, line=LINEC, line_w=Pt(0.75), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, Inches(0.85), Inches(4.7), Inches(11.6), Inches(0.35),
    [[("FORMAL COMPATIBILITY RULES  (modeled as SysML v2 constraints; enforced by the configuration sweep)", 11, True, NAVY)]])
bullets(s, Inches(0.85), Inches(5.12), Inches(5.9), Inches(1.4), [
    bullet_para(lead="P1  Cell-count.  ", size=12, text="Battery cells must lie within airframe min/max (e.g. no 4S pack on a 6S-only frame)."),
    bullet_para(lead="V1–V4  Video chain.  ", size=12, text="Camera output format must match VTX / ground-receiver format end-to-end."),
], size=12)
bullets(s, Inches(6.95), Inches(5.12), Inches(5.6), Inches(1.4), [
    bullet_para(lead="R1–R2  RF band.  ", size=12, text="Air/ground radio bands must match per link."),
    bullet_para(lead="Connector.  ", size=12, text="Whole power chain standardized on XT60 (battery → filter → airframe)."),
], size=12)
txt(s, Inches(6.9), Inches(3.7), Inches(5.8), Inches(0.6),
    [[("* Phase 4 downlink is a deferred future capability, shown for completeness; not part of the committed Phase 1–3 system.", 9.5, False, MUTED, True)]])

# ---- Standards Profile (StdV-1) ----
s = content_slide("Section 2 · Architecture", "Standards Profile  (StdV-1)")
txt(s, ML, Inches(1.28), CW, Inches(0.3),
    [[("The system is built on mature, open / de-facto-standard protocols — the interoperability basis behind the interfaces on the previous slide.", 12, False, MUTED)]])
stdv = [["Domain", "Standard / protocol", "Where used"]]
for a,b,c in [
    ("Flight control / autonomy","ArduPilot ArduCopter ≥ 4.5 (GPLv3)","FC firmware — AUTO/GUIDED, FS_/BATT_ failsafes"),
    ("Command & telemetry","MAVLink 2","FC ↔ SBC (UART) · FC ↔ GCS"),
    ("RC control link","ExpressLRS / CRSF","2.4 GHz control uplink + telemetry"),
    ("Ground-station app","QGroundControl 4.4+ (Apache-2.0)","Mission plan, telemetry, video, params"),
    ("Thermal video","USB Video Class (UVC 1.1)","Thermal camera → SBC"),
    ("FPV video","CVBS (analog composite)","FPV camera → VTX → ground VRX"),
    ("Navigation","u-blox UBX / NMEA 0183","GPS → flight controller"),
    ("AI inference runtime","RKNN (Rockchip NPU)","INT8 model on the NanoPi M5 (6 TOPS)"),
    ("Detection sizing","Johnson criteria","Thermal detect / recognize analysis"),
    ("Power connector","XT60","Battery → anti-spark → airframe"),
    ("Digital downlink (Phase 4)","OpenHD / WFB-ng over 802.11 monitor mode","Deferred thermal video downlink"),
    ("Regulatory","FAA Part 107 (+ BVLOS authorization)","Commercial UAS operation"),
]:
    stdv.append([a,b,c])
add_table(s, ML, Inches(1.68), CW, len(stdv),
          [Inches(2.85), Inches(4.6), Inches(4.64)], stdv,
          aligns=[PP_ALIGN.LEFT,PP_ALIGN.LEFT,PP_ALIGN.LEFT], row_h=Inches(0.37), body_size=10, header_size=10.5)
txt(s, ML, SH-Inches(0.55), CW, Inches(0.3),
    [[("StdV-1 (DoDAF Standards Profile) — no proprietary lock-in on the committed build; every committed interface rides an open or de-facto-standard protocol.", 9.5, False, MUTED, True)]])

# ============================================================================
# SECTION 3 — BEHAVIOR
# ============================================================================
section_divider("3", "Functional & Behavioral Design",
    "The functional architecture: a flight-mode state machine with failsafes, and the autonomous "
    "detect–investigate–classify loop that runs on the onboard computer during a sortie.")

# ---- 19. State machine ----
s = content_slide("Section 3 · Behavior", "Flight-Mode State Machine & Failsafes")
framed_image(s, im("drone_state_diagram.png"), ML, Inches(1.4), Inches(7.3), Inches(5.15))
caption(s, ML, Inches(6.6), Inches(7.3), "FlightMode state machine — disarmed · armed · flying · returnToLaunch · land")
bullets(s, Inches(8.15), Inches(1.5), Inches(4.55), Inches(5.0), [
    bullet_para(lead="Nominal path.  ", size=13.5, text="disarmed → armed → flying (route execution) → returnToLaunch → land (touchdown & disarm)."),
    bullet_para("Failsafe transitions (safety design):", size=13.5, bold=True, color=NAVY, bullet=True),
    bullet_para(lead="Link loss.  ", size=12.5, level=1, text="flying → RTL when ELRS is lost past the FC failsafe timeout (R7_BHV_LINKLOSS_RTL)."),
    bullet_para(lead="Low battery.  ", size=12.5, level=1, text="→ RTL when usable energy hits the return reserve (R6_BHV_RTL_RESERVE)."),
    bullet_para(lead="Realization.  ", size=13.5, text="Transitions map to ArduPilot FS_* / BATT_* failsafe parameters; executeRtl / touchDownAndDisarm are modeled actions."),
    bullet_para(lead="Open policy.  ", size=12.5, text="RTL reserve-energy threshold value is TBD (tracked as an open item)."),
])

# ---- 20. Autonomy loop ----
s = content_slide("Section 3 · Behavior", "Autonomous Detect · Investigate · Classify Loop")
framed_image(s, im("sweep_and_detect_action_flow.png"), ML, Inches(1.35), Inches(3.35), Inches(5.05))
caption(s, ML, Inches(6.45), Inches(3.35), "SweepAndDetect (UC-5a)")
framed_image(s, im("investigate_and_classify_action_flow.png"), Inches(4.15), Inches(1.35), Inches(3.9), Inches(5.05))
caption(s, Inches(4.15), Inches(6.45), Inches(3.9), "InvestigateAndClassify (UC-5b)")
rect(s, Inches(8.35), Inches(1.4), Inches(4.35), Inches(5.0), PANEL, line=LINEC, line_w=Pt(0.75), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, Inches(8.55), Inches(1.55), Inches(4.0), Inches(0.35), [[("HOW THE LOOP RUNS", 11, True, ACCENT)]])
bullets(s, Inches(8.55), Inches(1.98), Inches(3.95), Inches(4.3), [
    bullet_para(lead="Detect.  ", size=12.5, text="Per-frame onboard inference; when detectionConfidence ≥ threshold → mark POI, alert operator, emit targetDetectedEvt."),
    bullet_para(lead="Investigate.  ", size=12.5, text="Reroute to target; classify in a loop until classificationConfidence ≥ threshold, else adjust orbit and retry."),
    bullet_para(lead="Resume.  ", size=12.5, text="Report classification, resume the survey route, emit investigationCompleteEvt."),
    bullet_para(lead="Seam.  ", size=12.5, text="Both halves formally send their events to missionContext — the receiver that ties functions to the system."),
    bullet_para(lead="Allocation.  ", size=12.5, text="Runs on the SBC (missionApp + rknnRuntime); routes commanded to the FC over MAVLink."),
])

# ============================================================================
# SECTION 4 — ANALYSES
# ============================================================================
section_divider("4", "Engineering Analyses",
    "Four quantitative analyses verify the design against the driving requirements — endurance, "
    "thermal detection, RF link budget, and cost / physical integration — each with its method and "
    "fidelity stated.")

# ---- 22. Design selection / reference build ----
s = content_slide("Section 4 · Analyses", "Detailed Design — Selected Configuration")
txt(s, ML, Inches(1.28), CW, Inches(0.3),
    [[("Components down-selected from ~110 real market candidates via trade studies; selections are version-controlled and machine-enforced in the analysis.", 12, False, MUTED)]])
sel = [["Role", "Selected component", "Basis"]]
for a,b,c in [
    ("Airframe","iFlight Chimera9 ECO 9\" (PNP)","Best endurance/$ of SBC-capable frames; SBC fits deck"),
    ("Battery (flight)","Upgrade Energy 6S 12 Ah Amprius","57 min hover; XT60; in-stock Amprius"),
    ("Thermal camera","PurpleRiver Mini 640 (640×512, 12 µm)","Meets Johnson detect/recognize; USB-UVC"),
    ("Onboard compute","NanoPi M5 (RK3576, 6 TOPS)","Mature RKNN toolchain; ≤10 W passive"),
    ("Control link","iFlight TD ELRS + TX12 MkII + Hermes","True-diversity 2.4 GHz; huge link margin"),
    ("Video (ground)","Skydroid UVC RX + TrueRC X-AIR patch","Closes 2.8 km analog link with margin"),
    ("FC firmware","ArduPilot ArduCopter ≥ 4.5","AUTO/GUIDED + FS_/BATT_ failsafes ($0)"),
    ("GCS app","QGroundControl 4.4+","All laptop functions; macOS-compatible ($0)"),
]:
    sel.append([a,b,c])
add_table(s, ML, Inches(1.75), CW, len(sel),
          [Inches(2.35), Inches(4.4), Inches(5.34)], sel,
          aligns=[PP_ALIGN.LEFT,PP_ALIGN.LEFT,PP_ALIGN.LEFT], row_h=Inches(0.44), body_size=10.5, header_size=11)
txt(s, ML, Inches(6.35), CW, Inches(0.5),
    [[("Reference build: ", 10.5, True, NAVY),
      ("Chimera9 ECO + 6S 12 Ah Amprius + Mini 640 + NanoPi M5 → ≈ 58 min hover, ≈ $1,850 committed system cost (≤ $2,500 R4), SBC fits the deck.", 10.5, False, TEXT)]])

# ---- Selected-hardware photo gallery ----
import math as _math
def hw_card(slide, l, t, w, h, img, name, spec):
    rect(slide, l, t, w, h, WHITE, line=LINEC, line_w=Pt(1), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt_h = Inches(0.74)
    img_box_h = h - txt_h - Inches(0.12)
    if img and os.path.exists(img):
        fl, ft, fw, fh = fit_box(img, l+Inches(0.12), t+Inches(0.1), w-Inches(0.24), img_box_h)
        slide.shapes.add_picture(img, fl, ft, fw, fh)
    else:
        rect(slide, l+Inches(0.2), t+Inches(0.2), w-Inches(0.4), img_box_h-Inches(0.15), PANEL,
             line=LINEC, line_w=Pt(0.5), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        txt(slide, l+Inches(0.2), t+img_box_h/2-Inches(0.1), w-Inches(0.4), Inches(0.4),
            [[("see spec kit", 10, False, MUTED, True)]], align=PP_ALIGN.CENTER)
    rect(slide, l+Inches(0.12), t+h-txt_h+Inches(0.02), w-Inches(0.24), Pt(1.2), LINEC)
    txt(slide, l+Inches(0.12), t+h-txt_h+Inches(0.07), w-Inches(0.24), Inches(0.34),
        [[(name, 12, True, NAVY)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(slide, l+Inches(0.12), t+h-txt_h+Inches(0.42), w-Inches(0.24), Inches(0.34),
        [[(spec, 9.5, False, MUTED)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)

def hardware_gallery(kicker, title, items, note=None):
    s = content_slide(kicker, title)
    n=len(items); cols=3; rows=_math.ceil(n/cols)
    gap=Inches(0.32)
    area_top=Inches(1.4); area_h = Inches(4.8) if note else Inches(5.15)
    card_w = (CW - gap*(cols-1))/cols
    card_h = min((area_h - gap*(rows-1))/rows, Inches(2.6))
    block_h = card_h*rows + gap*(rows-1)
    top = area_top + (area_h - block_h)//2
    for i,(img,name,spec) in enumerate(items):
        c=i%cols; r=i//cols
        l = ML + c*(card_w+gap); t = top + r*(card_h+gap)
        hw_card(s, l, t, card_w, card_h, img, name, spec)
    if note:
        rect(s, ML, Inches(6.5), CW, Inches(0.6), PANEL, line=LINEC, line_w=Pt(0.75), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        txt(s, ML+Inches(0.2), Inches(6.5), CW-Inches(0.4), Inches(0.6), note, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)
    return s

hardware_gallery("Section 4 · Analyses  ·  Detailed Design", "Selected Hardware — Air Vehicle & Payload", [
    (comp("airframe.png"),  "iFlight Chimera9 ECO 9\"", "PNP airframe · 729 g · 6S · XING-E 2809"),
    (comp("battery.png"),   "Upgrade Energy 6S 12 Ah",  "Amprius Si-anode Li-ion · XT60 · ~57 min"),
    (comp("thermal.png"),   "PurpleRiver Mini 640",     "640×512 · 12 µm · 13 mm · USB-UVC"),
    (comp("sbc.png"),       "NanoPi M5 (RK3576)",       "6 TOPS NPU · ≤10 W passive · 90×62 mm"),
    (comp("receiver.png"),  "iFlight True-Diversity ELRS","2.4 GHz · dual-antenna · 250 mW"),
    (comp("antispark.png"), "iFlight Anti-Spark Filter","XT60 · 6S · 150 A · inrush limiter"),
])

hardware_gallery("Section 4 · Analyses  ·  Detailed Design", "Selected Hardware — Ground Control & Support", [
    (comp("radio.png"),   "RadioMaster TX12 Mark II", "ELRS 2.4 GHz · EdgeTX · hall gimbals"),
    (comp("antenna.png"), "TrueRC X-AIR 5.8 MK II",   "10 dBic patch · 120° · RHCP (ground)"),
    (comp("charger.png"), "HOTA D6 Pro",              "AC 200 / DC 650 W · dual · Li-ion + LiPo"),
], note=[[("Also in the ground / support kit:  ", 11, True, NAVY),
    ("Skydroid 150CH 5.8 GHz UVC video receiver ($44)  ·  HGLRC Hermes ELRS USB telemetry dongle ($16)  ·  GPS (iFlight BLITZ M10, bundled with the airframe)  ·  MacBook Air GCS (existing).", 11, False, TEXT)]])

# ---- 23. Endurance analysis ----
s = content_slide("Section 4 · Analyses", "Analysis 1 — Endurance (Flight Time)")
rect(s, ML, Inches(1.3), Inches(6.05), Inches(1.55), PANEL, line=LINEC, line_w=Pt(0.75), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, Inches(0.8), Inches(1.42), Inches(5.6), Inches(0.3), [[("METHOD & FIDELITY", 11, True, ACCENT)]])
txt(s, Inches(0.8), Inches(1.75), Inches(5.6), Inches(1.05),
    [[("Momentum / actuator-disk (blade-element) model with a forward-flight parasitic-drag term — the same physics family as eCalc. Holistic sweep crosses airframe × battery × payload, filtered by the compatibility rules. ", 10.5, False, TEXT),
      ("Fidelity: first-order comparative; FoM 0.65, η 0.80, ρ 1.225, C_d 1.0. Not yet validated against measured hover current.", 10.5, True, RED)]], line_spacing=1.03)
bullets(s, ML, Inches(3.05), Inches(6.05), Inches(3.5), [
    bullet_para(lead="Requirements verified.  ", size=13, text="R6 (≥ 30 min) and R8 stretch (≥ 60 min); R4 cost checked per config."),
    bullet_para(lead="Scope.  ", size=13, text="34 compatible configs (airframe × battery × VTX crossed; T13 + NanoPi M5 fixed) — filtered by the cell-count rule (P1)."),
    bullet_para(lead="Result — baseline.  ", size=13, text="Chimera9 ECO + 6S 12 Ah Amprius → 58.4 min hover at ~15% throttle (16.85 W payload), 1,751 g AUW."),
    bullet_para(lead="Power bucket.  ", size=13, text="Slow cruise (2.23 m/s) can exceed hover via translational lift — cruise ≈ 60 min, modeled explicitly."),
])
# result stats
stat(s, Inches(7.0), Inches(1.3), Inches(2.75), Inches(1.5), "58.4 min", "Baseline hover endurance", accent=OKGRN)
stat(s, Inches(9.95), Inches(1.3), Inches(2.75), Inches(1.5), "≈ 15 %", "Hover throttle (large margin)")
stat(s, Inches(7.0), Inches(3.0), Inches(2.75), Inches(1.5), "R6 ✓ / R8 ✗", "30 min met; 60 min stretch not", accent=AMBER)
stat(s, Inches(9.95), Inches(3.0), Inches(2.75), Inches(1.5), "34", "Compatible configs swept")
rect(s, Inches(7.0), Inches(4.75), Inches(5.7), Inches(1.75), PANEL, line=LINEC, line_w=Pt(0.75), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, Inches(7.2), Inches(4.88), Inches(5.4), Inches(0.3), [[("INTEGRATED WITH THE MODEL", 11, True, ACCENT)]])
txt(s, Inches(7.2), Inches(5.2), Inches(5.4), Inches(1.2),
    [[("The script reads component + battery data directly from candidates.sysml, applies the model's compatibility rules, and writes back ranked results (CSV + a SysML v2 instance table) — a closed loop, not a standalone spreadsheet.", 10.5, False, TEXT)]], line_spacing=1.05)

# ---- 24. Cost vs flight-time scatter ----
s = content_slide("Section 4 · Analyses", "Endurance vs. Cost — Full Configuration Trade Space")
framed_image(s, COSTCHART, ML, Inches(1.35), Inches(8.4), Inches(5.3))
caption(s, ML, Inches(6.62), Inches(8.4), "Every compatible real configuration: total system cost vs. flight time, colored by airframe (R4/R6/R8 reference lines)")
bullets(s, Inches(9.05), Inches(1.5), Inches(3.65), Inches(5.0), [
    bullet_para(lead="Reading it.  ", size=13, text="Each point is a real, purchasable build. Vertical line = $2,500 (R4); horizontals = 30 min (R6) and 60 min (R8)."),
    bullet_para(lead="Result.  ", size=13, text="The entire selected trade space sits left of the R4 cap and above R6 — cost is not the binding constraint."),
    bullet_para(lead="Selected region.  ", size=13, text="The reference build lands at ≈ $1,850 / 58 min — comfortable margin on both axes."),
    bullet_para(lead="Stretch (R8).  ", size=13, text="60 min hover is not reached by any config with the fixed T13 + NanoPi M5 payload; cruise endurance approaches it."),
])

# ---- 25. Thermal detection ----
s = content_slide("Section 4 · Analyses", "Analysis 2 — Thermal Detection (Johnson Criteria)")
rect(s, ML, Inches(1.3), Inches(6.05), Inches(1.55), PANEL, line=LINEC, line_w=Pt(0.75), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, Inches(0.8), Inches(1.42), Inches(5.6), Inches(0.3), [[("METHOD & FIDELITY", 11, True, ACCENT)]])
txt(s, Inches(0.8), Inches(1.75), Inches(5.6), Inches(1.05),
    [[("Geometric optics: ground sample distance (GSD) from sensor pitch, lens, altitude → pixels-on-target vs. Johnson criteria (detect ≥ 1.5 px, recognize ≥ 4 px across the target's minimum dimension). ", 10.5, False, TEXT),
      ("Fidelity: deterministic; assumes ≥ 5 °C target-to-background thermal contrast (daytime IR signature). Species-level classification ultimately relies on the trained AI model (Phase 3) — field capture pending.", 10.5, True, RED)]], line_spacing=1.02)
# T13 results table
th = [["Target @ altitude", "Min px", "Verdict"]]
for a,b,c in [
    ("0.5 m target @ 90 m","6.0","Recognize ✓"),
    ("0.5 m target @ 120 m","4.5","Recognize ✓"),
    ("Deer @ 120 m (R3.1)","4.5","Detect ✓"),
    ("Deer @ 90 m","6.0","Detect / recog."),
    ("Human @ 90 m","6.0","Detect ✓"),
    ("Turkey @ 120 m","2.7","Marginal"),
]:
    th.append([a,b,c])
bcth = {}
for i in range(1,len(th)):
    v = th[i][2]
    bcth[(i,2)] = OKGRN if "✓" in v else (AMBER if "Marginal" in v or "recog" in v else TEXT)
add_table(s, ML, Inches(3.05), Inches(6.05), len(th),
          [Inches(3.35), Inches(1.1), Inches(1.6)], th,
          aligns=[PP_ALIGN.LEFT,PP_ALIGN.CENTER,PP_ALIGN.CENTER], row_h=Inches(0.38), body_size=10.5, header_size=10.5, body_colors=bcth)
# right column
stat(s, Inches(7.0), Inches(1.3), Inches(2.75), Inches(1.5), "8.3 cm", "GSD per pixel @ 90 m")
stat(s, Inches(9.95), Inches(1.3), Inches(2.75), Inches(1.5), "11.1 cm", "GSD per pixel @ 120 m")
stat(s, Inches(7.0), Inches(3.0), Inches(2.75), Inches(1.5), "640×512", "Sensor, 12 µm pitch, 13 mm lens", accent=OKGRN)
stat(s, Inches(9.95), Inches(3.0), Inches(2.75), Inches(1.5), "R3.1 ✓", "Deer detection @ 120 m", accent=OKGRN)
rect(s, Inches(7.0), Inches(4.75), Inches(5.7), Inches(1.75), PANEL, line=LINEC, line_w=Pt(0.75), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, Inches(7.2), Inches(4.88), Inches(5.4), Inches(0.3), [[("INTERPRETATION", 11, True, ACCENT)]])
txt(s, Inches(7.2), Inches(5.2), Inches(5.4), Inches(1.2),
    [[("Optics deliver detection-grade pixels-on-target for deer/human across 90–120 m and recognition-grade for a 0.5 m target. Robust species classification (R3.2) is the AI model's job and the key validation risk — confirmed by field thermal capture, not geometry alone.", 10.5, False, TEXT)]], line_spacing=1.03)

# ---- 26. RF link budget ----
s = content_slide("Section 4 · Analyses", "Analysis 3 — RF Link Budget (2.8 km / R7)")
txt(s, ML, Inches(1.28), CW, Inches(0.55),
    [[("Method & fidelity:  ", 11.5, True, ACCENT),
      ("free-space path loss (FSPL) at each link's frequency, design target ≥ 10 dB fade margin at 2.8 km. Fidelity: free-space model; the fade margin absorbs foliage / multipath / polarization loss. Not yet range-tested.", 11.5, False, TEXT)]], line_spacing=1.03)
rf = [["RF link", "TX", "Margin @ 2.8 km", "Reliable range", "Verdict"]]
for a,b,c,d2,e in [
    ("Analog FPV video (5.8 GHz + patch)","32 dBm","+15.3 dB","5.2 km","PASS"),
    ("Control uplink (2.4 GHz ELRS)","24 dBm","+23.0 dB","12.5 km","PASS"),
    ("Telemetry downlink (2.4 GHz ELRS)","20 dBm","+19.0 dB","7.9 km","PASS"),
    ("OpenHD thermal video (5.8 GHz)*","29 dBm","+11.3 dB","3.3 km","PASS"),
]:
    rf.append([a,b,c,d2,e])
bcrf = {(i,4): OKGRN for i in range(1,len(rf))}
for i in range(1,len(rf)): bcrf[(i,2)] = OKGRN
add_table(s, ML, Inches(2.05), CW, len(rf),
          [Inches(5.0), Inches(1.4), Inches(2.4), Inches(1.9), Inches(1.39)], rf,
          aligns=[PP_ALIGN.LEFT,PP_ALIGN.CENTER,PP_ALIGN.CENTER,PP_ALIGN.CENTER,PP_ALIGN.CENTER],
          row_h=Inches(0.46), body_size=11, header_size=10.5, body_colors=bcrf)
bullets(s, ML, Inches(4.55), Inches(7.4), Inches(2.0), [
    bullet_para(lead="All links pass 2.8 km with ≥ 10 dB margin.  ", size=13, text="Control/telemetry (ELRS) close the range many times over and are never limiting."),
    bullet_para(lead="Binding link.  ", size=13, text="Analog FPV video — the TrueRC X-AIR patch antenna (10 dBic) is the enabling component; the stock omni would fail the 10 dB target."),
    bullet_para(text="* Phase 4 (deferred). OpenHD thermal downlink shown for completeness; dual-diversity Foxeer panels give +11.3 dB.", size=12.5, color=MUTED),
])
stat(s, Inches(8.3), Inches(4.55), Inches(2.1), Inches(1.9), "4 / 4", "RF links pass R7", accent=OKGRN)
stat(s, Inches(10.6), Inches(4.55), Inches(2.1), Inches(1.9), "+10 dB", "Fade-margin design target")
txt(s, ML, Inches(6.68), Inches(7.4), Inches(0.42),
    [[("Fade margin = ", 9.5, True, NAVY),
      ("dB of received-signal headroom above the reliable-link minimum. The ≥ 10 dB target is a ~10× power cushion so the link holds through momentary signal fades (multipath, foliage, bank angle).", 9.5, False, MUTED)]], line_spacing=1.0)

# ---- 27. Physical integration ----
s = content_slide("Section 4 · Analyses", "Analysis 4 — Physical Integration (“Does It Fit?”)")
rect(s, ML, Inches(1.3), Inches(6.05), Inches(1.75), PANEL, line=LINEC, line_w=Pt(0.75), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, Inches(0.8), Inches(1.42), Inches(5.6), Inches(0.3), [[("METHOD & FIDELITY", 11, True, ACCENT)]])
txt(s, Inches(0.8), Inches(1.75), Inches(5.6), Inches(1.25),
    [[("Dimensional model: parses each airframe's usable deck footprint and each battery / SBC envelope, then emits a 3-tier fit verdict (fits / marginal / no-fit, ±12 mm tolerance) with deck-margin. ", 10.5, False, TEXT),
      ("Fidelity: deck dimensions are mostly ESTIMATES (≈ 0.27 × wheelbase); battery envelopes validated against the one confirmed pack. Confirm the selected frame's real deck before build.", 10.5, True, RED)]], line_spacing=1.0)
bullets(s, ML, Inches(3.25), Inches(6.05), Inches(3.3), [
    bullet_para(lead="Binding constraint.  ", size=13, text="The NanoPi M5 footprint (90 × 62 mm) — the thermal camera is nose-mounted and never an issue."),
    bullet_para(lead="Selected frame.  ", size=13, text="Chimera9 ECO deck ≈ 110 × 70 mm → SBC fits with ~8 mm spare (a key reason it was chosen over the endurance leader)."),
    bullet_para(lead="General pattern.  ", size=13, text="9–10\" frames fit cleanly; 7.5–9\" are marginal (custom deck); the pure endurance winner (DarwinFPV 129) is NO-FIT."),
    bullet_para(lead="Status.  ", size=13, text="Informational — configs are flagged, never dropped. Fit drove the airframe lock."),
])
stat(s, Inches(7.0), Inches(1.3), Inches(2.75), Inches(1.6), "90×62 mm", "SBC footprint (binding)", accent=AMBER)
stat(s, Inches(9.95), Inches(1.3), Inches(2.75), Inches(1.6), "~8 mm", "Deck spare on selected frame", accent=OKGRN)
stat(s, Inches(7.0), Inches(3.1), Inches(2.75), Inches(1.6), "Fits ✓", "Chimera9 ECO (110×70 deck)", accent=OKGRN)
stat(s, Inches(9.95), Inches(3.1), Inches(2.75), Inches(1.6), "3D-print", "SBC mount + 30 mm fan (Phase 2)")
rect(s, Inches(7.0), Inches(4.85), Inches(5.7), Inches(1.65), PANEL, line=LINEC, line_w=Pt(0.75), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, Inches(7.2), Inches(4.98), Inches(5.4), Inches(0.3), [[("WHY IT MATTERS", 11, True, ACCENT)]])
txt(s, Inches(7.2), Inches(5.3), Inches(5.4), Inches(1.1),
    [[("Endurance and fit trade against each other: the lightest/longest-flying frames are too small to carry the compute. The selected frame is the best endurance-per-dollar option that physically integrates the payload.", 10.5, False, TEXT)]], line_spacing=1.03)

# ---- 28. Cost / budget ----
s = content_slide("Section 4 · Analyses", "Analysis 5 — Cost & Budget (R4 ≤ $2,500)")
txt(s, ML, Inches(1.28), CW, Inches(0.3),
    [[("Method:  ", 11.5, True, ACCENT),
      ("phased bill of materials rolled up from candidates.sysml with actual quoted prices for locked parts; laptop (existing) and reusable support gear (dev packs, charger) excluded from the R4 integrated-system figure.", 11.5, False, TEXT)]], line_spacing=1.03)
cost = [["Build phase", "Scope", "Subtotal"]]
for a,b,c in [
    ("Phase 1","Flight + FPV downlink + waypoints","$1,371"),
    ("Phase 2","Thermal camera + onboard SBC","$809"),
    ("Phase 3","AI detection + autonomy (software only)","$0"),
    ("Committed (1–3)","As-flown system","$2,180"),
    ("R4 integrated","Excl. reusable dev packs + charger","≈ $1,848"),
    ("Phase 4 (deferred)","OpenHD digital downlink","$159"),
]:
    cost.append([a,b,c])
bcc = {(4,0):NAVY,(5,0):OKGRN}
add_table(s, ML, Inches(1.85), Inches(7.3), len(cost),
          [Inches(2.2), Inches(3.7), Inches(1.4)], cost,
          aligns=[PP_ALIGN.LEFT,PP_ALIGN.LEFT,PP_ALIGN.RIGHT], row_h=Inches(0.42), body_size=11, header_size=11)
# margin visual
stat(s, Inches(8.15), Inches(1.85), Inches(4.55), Inches(1.35), "≈ $1,848", "Committed system cost (R4 basis)", accent=OKGRN)
stat(s, Inches(8.15), Inches(3.35), Inches(4.55), Inches(1.35), "≈ 26 %", "Margin under the $2,500 R4 cap", accent=OKGRN)
rect(s, Inches(8.15), Inches(4.85), Inches(4.55), Inches(1.65), PANEL, line=LINEC, line_w=Pt(0.75), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, Inches(8.35), Inches(4.98), Inches(4.2), Inches(0.3), [[("R4 VERDICT", 11, True, ACCENT)]])
txt(s, Inches(8.35), Inches(5.3), Inches(4.2), Inches(1.1),
    [[("PASS. ", 11, True, OKGRN),
      ("Committed system (≈ $1,848) and all four phases (≈ $2,007) are both under $2,500. Thermal camera ($650) is the single largest line.", 10.5, False, TEXT)]], line_spacing=1.03)
txt(s, ML, Inches(6.45), Inches(7.2), Inches(0.4),
    [[("The endurance sweep reports ≈ $1,674 using cost-representative radio/RX; the BOM figure uses the actual selected parts, hence higher.", 9.5, False, MUTED, True)]])

# ---- Systems Measures / performance scorecard (SV-7) ----
s = content_slide("Section 4 · Analyses", "Performance Measures  (SV-7)")
txt(s, ML, Inches(1.28), CW, Inches(0.3),
    [[("Key performance parameters vs. their requirement thresholds, with the as-designed (analysis) value and margin — the quantitative scorecard behind the analyses.", 12, False, MUTED)]])
sv7 = [["Measure", "Requirement", "As-designed (predicted)", "Margin", "Status"]]
for r_ in [
    ("Hover endurance","R6 ≥ 30 min","58.4 min (momentum model)","+95%","PASS"),
    ("Stretch endurance","R8 ≥ 60 min","58.4 min hover (~60 cruise)","−3%","NOT MET"),
    ("Surveillance range","R7 ≥ 2.8 km","≥ 3.3 km (weakest RF link)","+0.5 km","PASS"),
    ("RF fade margin","Design ≥ 10 dB @ 2.8 km","+11.3 dB (weakest link)","+1.3 dB","PASS"),
    ("Thermal detection @ 120 m","R3.1 ≥ 1.5 px","4.5 px on-target (deer)","+3.0 px","PASS"),
    ("Thermal recognition @ 90 m","R3_CAM_RES ≥ 4 px","6.0 px (0.5 m target)","+2.0 px","PASS"),
    ("System cost","R4 ≤ $2,500","≈ $1,848 committed","+$652 (26%)","PASS"),
    ("SBC power","R4_SBC_PWR ≤ 10 W","≤ 10 W passive (RK3576)","at limit","PASS"),
    ("Payload compute fit","Deck ≥ SBC footprint","110×70 vs 90×62 mm","+8 mm","PASS*"),
    ("Cruise speed","R2 = 2.23 m/s","design point (control-tuned)","—","PLANNED"),
]:
    sv7.append(list(r_))
bcsv = {}
for i in range(1,len(sv7)):
    st = sv7[i][4]; mg = sv7[i][3]
    bcsv[(i,4)] = OKGRN if st.startswith("PASS") else (RED if "NOT MET" in st else SLATE)
    bcsv[(i,3)] = RED if mg=="−3%" else (AMBER if mg in ("at limit","—") else OKGRN)
add_table(s, ML, Inches(1.68), CW, len(sv7),
          [Inches(2.95), Inches(2.35), Inches(3.4), Inches(1.6), Inches(1.79)], sv7,
          aligns=[PP_ALIGN.LEFT,PP_ALIGN.LEFT,PP_ALIGN.LEFT,PP_ALIGN.CENTER,PP_ALIGN.CENTER],
          row_h=Inches(0.4), body_size=10, header_size=10.5, body_colors=bcsv)
txt(s, ML, SH-Inches(0.55), CW, Inches(0.3),
    [[("SV-7 (DoDAF Systems Measures). All hard requirements met with margin; only the R8 stretch endurance is not met (documented). *Compute fit is marginal (custom deck); speed is a control-tuned target verified by flight test.", 9.5, False, MUTED, True)]])

# ============================================================================
# SECTION 5 — V&V, RISK, CLOSURE
# ============================================================================
section_divider("5", "Verification, Risk & Closure",
    "How each requirement will be verified, the manufacturing approach, the risk posture, remaining "
    "open items, and the CDR exit assessment.")

# ---- 30. VCRM ----
s = content_slide("Section 5 · Closure", "Verification Cross-Reference Matrix (VCRM)")
txt(s, ML, Inches(1.28), CW, Inches(0.3),
    [[("At CDR, design compliance is shown by analysis; demonstration/test are planned and gated at the Test Readiness Review.", 12, False, MUTED)]])
vc = [["Req", "Requirement", "Method", "Verification approach", "Status @ CDR"]]
for rid, desc, meth, appr, stat_ in [
    ("R1","90–120 m AGL altitude hold","A, T","Control-loop analysis; logged flight test","Analysis ✓ · Test planned"),
    ("R2","2.23 m/s ground speed","A, T","Analysis; logged flight test","Analysis ✓ · Test planned"),
    ("R3.1","Detect target @ 120 m","A, D","Johnson-criteria analysis; field thermal demo","Analysis ✓ · Demo planned"),
    ("R3.2","Classify species @ 90 m","A, D, T","AI-model eval + field capture","Analysis ✓ · AI test planned"),
    ("R4","System cost ≤ $2,500","I","BOM inspection","Verified ✓"),
    ("R5","Minimize DIY soldering","I","Design inspection (PNP / pre-soldered)","Verified ✓"),
    ("R6","≥ 30 min endurance","A, T","Endurance model; measured hover","Analysis ✓ · Test planned"),
    ("R7","2.8 km range in wind","A, T","FSPL link budget; range test","Analysis ✓ · Test planned"),
    ("R8","60 min stretch endurance","A","Endurance model","Not met (documented)"),
]:
    vc.append([rid, desc, meth, appr, stat_])
bcv = {}
for i in range(1,len(vc)):
    st = vc[i][4]
    bcv[(i,4)] = OKGRN if st.startswith("Verified") else (RED if "Not met" in st else SLATE)
add_table(s, ML, Inches(1.72), CW, len(vc),
          [Inches(0.9), Inches(3.0), Inches(1.15), Inches(4.0), Inches(3.04)], vc,
          aligns=[PP_ALIGN.CENTER,PP_ALIGN.LEFT,PP_ALIGN.CENTER,PP_ALIGN.LEFT,PP_ALIGN.LEFT],
          row_h=Inches(0.44), body_size=10, header_size=10.5, body_colors=bcv)

# ---- 31. Manufacturing ----
s = content_slide("Section 5 · Closure", "Manufacturing & Producibility")
bullets(s, ML, Inches(1.4), Inches(6.1), Inches(5.1), [
    bullet_para(lead="Minimize soldering (R5).  ", size=13.5, text="PNP airframe ships with pre-soldered power distribution, plug-in connectors, and screw-terminal motors. The only solder joint in the committed build is the SBC UBEC's 12 V leads."),
    bullet_para(lead="COTS-first.  ", size=13.5, text="~110 components traded were all real, purchasable parts; the selected BOM is off-the-shelf except the fabricated SBC mount."),
    bullet_para(lead="Fabricated parts.  ", size=13.5, text="3D-printed SBC deck + 30 mm fan + heat-set hardware (free-CAD pathway documented); ≈ $15."),
    bullet_para(lead="Standardization.  ", size=13.5, text="Whole power chain on XT60; single control ecosystem (ELRS) end-to-end."),
    bullet_para(lead="Phased integration.  ", size=13.5, text="Three incremental builds, each independently testable before the next is added."),
])
rect(s, Inches(7.05), Inches(1.4), Inches(5.65), Inches(5.1), PANEL, line=LINEC, line_w=Pt(0.75), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, Inches(7.3), Inches(1.58), Inches(5.2), Inches(0.35), [[("PHASED BUILD PLAN", 11.5, True, ACCENT)]])
bullets(s, Inches(7.3), Inches(2.05), Inches(5.15), Inches(4.3), [
    bullet_para(lead="Phase 1.  ", size=13, text="Airframe + ELRS + battery + FPV/GPS → LOS manual flight, waypoint routes, video downlink."),
    bullet_para(lead="Phase 2.  ", size=13, text="Thermal camera + NanoPi M5 → live onboard inference feed (no recording, no downlink)."),
    bullet_para(lead="Phase 3.  ", size=13, text="AI detection + MAVLink autonomous route modification (software only)."),
    bullet_para(lead="Phase 4 (deferred).  ", size=13, text="OpenHD digital thermal downlink to the ground station — future capability.", color=MUTED),
])

# ---- Project Timeline (PV-2) ----
s = content_slide("Section 5 · Closure", "Project Timeline  (PV-2)")
txt(s, ML, Inches(1.26), CW, Inches(0.3),
    [[("The build roadmap: three committed phases (1–3) plus a deferred Phase 4. Detailed design is complete; the next gate is procurement.", 12, False, MUTED)]])
rect(s, ML, Inches(1.72), CW, Inches(0.5), NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, ML+Inches(0.22), Inches(1.72), CW-Inches(0.44), Inches(0.5),
    [[("▸ NOW — CDR: detailed design complete, verified by analysis.   Next gate → procurement, then phased build & test.", 12, True, WHITE)]], anchor=MSO_ANCHOR.MIDDLE)
_phases = [
    ("PHASE 1","Committed","Basic flight + FPV + waypoints",
     ["Airframe · ELRS · battery","FPV cam · GPS · patch antenna","ArduPilot + QGroundControl"],
     "LOS manual flight, waypoint routes, video to laptop","≈ $1,371"),
    ("PHASE 2","Committed","Thermal payload + onboard compute",
     ["PurpleRiver Mini 640 (thermal)","NanoPi M5 · UBEC · 3D-print mount"],
     "Live onboard thermal inference (no recording / downlink)","≈ $809"),
    ("PHASE 3","Committed","AI detection + autonomy (software)",
     ["RKNN model deploy (D-1)","Mission app + MAVLink route mod (D-2)"],
     "Real-time inference drives autonomous re-route","$0 HW"),
    ("PHASE 4","Deferred","OpenHD digital thermal downlink",
     ["Air / ground Wi-Fi adapters","Diversity antennas · ground VM"],
     "Live thermal / AI video to the ground station","≈ $159 (future)"),
]
cw4 = (CW - Inches(0.3)*3)//4
ptop = Inches(2.48); ch4 = Inches(4.0)
for i,(ph,tag,title,items,deliv,cost) in enumerate(_phases):
    x = ML + i*(cw4 + Inches(0.3)); deferred = (tag=="Deferred")
    linecol = MUTED if deferred else ACCENT
    rect(s, x, ptop, cw4, ch4, WHITE, line=linecol, line_w=Pt(1.0 if deferred else 1.5), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, x, ptop, cw4, Inches(0.64), RGBColor(0x9A,0xA7,0xB2) if deferred else ACCENT, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, x, ptop+Inches(0.06), cw4, Inches(0.26), [[(ph, 13, True, WHITE)]], align=PP_ALIGN.CENTER)
    txt(s, x, ptop+Inches(0.36), cw4, Inches(0.22), [[(tag, 9.5, True, WHITE)]], align=PP_ALIGN.CENTER)
    txt(s, x+Inches(0.16), ptop+Inches(0.76), cw4-Inches(0.32), Inches(0.6), [[(title, 11.5, True, NAVY)]])
    yy = ptop+Inches(1.5)
    for it in items:
        txt(s, x+Inches(0.16), yy, cw4-Inches(0.32), Inches(0.42), [[("· ", 10, True, linecol), (it, 10, False, TEXT)]])
        yy += Inches(0.42)
    txt(s, x+Inches(0.16), ptop+ch4-Inches(1.2), cw4-Inches(0.32), Inches(0.8),
        [[("→ ", 10, True, linecol), (deliv, 10, False, MUTED, True)]], line_spacing=1.0)
    txt(s, x+Inches(0.16), ptop+ch4-Inches(0.38), cw4-Inches(0.32), Inches(0.3), [[(cost, 11.5, True, NAVY if not deferred else MUTED)]])
    if i < 3:
        txt(s, x+cw4-Inches(0.03), ptop+ch4//2-Inches(0.18), Inches(0.36), Inches(0.34), [[("▸", 15, True, linecol)]], align=PP_ALIGN.CENTER)
txt(s, ML, SH-Inches(0.5), CW, Inches(0.3),
    [[("PV-2 (DoDAF Project Timeline). Phases 1–3 = committed system; Phase 4 (grey) is deferred. Sequence and dependency, not calendar dates.", 9.5, False, MUTED, True)]])

# ---- 32. Risk ----
s = content_slide("Section 5 · Closure", "Risk Assessment")
rk = [["#", "Risk", "Sev.", "Mitigation"]]
for n, risk, sev, mit in [
    ("1","Design is modeled, not yet flown — headline numbers are predictions","High","Build Phase 1; instrumented hover to correlate endurance model within ±15%"),
    ("2","Species classification (R3.2) unproven — needs real thermal signatures + trained AI","High","Field thermal capture @ 90–120 m; train & validate model before claiming R3.2"),
    ("3","Regulatory — 2.8 km range is beyond visual line of sight (BVLOS); commercial use needs Part 107","Med","Operate within VLOS or obtain a BVLOS waiver / fly under Part 108; hold a Part 107 certificate; confirm current FAA rules"),
    ("4","Flight-battery single-source (Amprius; prior SKU already sold out)","Med","Qualify GNB 12 Ah fallback; document endurance delta"),
    ("5","Deck-fit uses estimated dimensions for the selected frame","Med","Measure real Chimera9 deck / vendor CAD before Phase 2 build"),
    ("6","Endurance model is first-order (FoM/η/C_d assumed)","Med","Correlate to measured hover current on first flights"),
]:
    rk.append([n, risk, sev, mit])
bcr = {}
for i in range(1,len(rk)):
    bcr[(i,2)] = RED if rk[i][2]=="High" else AMBER
add_table(s, ML, Inches(1.4), CW, len(rk),
          [Inches(0.5), Inches(5.2), Inches(0.95), Inches(5.44)], rk,
          aligns=[PP_ALIGN.CENTER,PP_ALIGN.LEFT,PP_ALIGN.CENTER,PP_ALIGN.LEFT],
          row_h=Inches(0.72), body_size=10.5, header_size=10.5, body_colors=bcr)
txt(s, ML, Inches(6.5), CW, Inches(0.4),
    [[("At CDR, “not yet built” is expected — the material risks are the estimates and the external (regulatory / classification) items, all with defined mitigations.", 10, False, MUTED, True)]])

# ---- 33. Open items ----
s = content_slide("Section 5 · Closure", "Open Items  (TBD / TBR / Action)")
oi = [["Item", "Type", "Owner action"]]
for a,b,c in [
    ("RTL reserve-energy threshold value (R6_BHV_RTL_RESERVE)","TBR","Set policy value before first autonomous flight"),
    ("FC firmware — reconciled in MODEL_ISSUES.md §B7 (RESOLVED)","RESOLVED","Issues log updated 2026-08-09 to match SELECTED_COMPONENTS"),
    ("Selected-frame deck & battery-bay dimensions are estimates","TBR","Confirm from vendor CAD / measurement (Phase 2 gate)"),
    ("SBC mount — 3D-printed deck not yet designed or modeled","Design","Design the mount + add it to model.sysml and CAD (Phase 2 gate)"),
    ("AF5 (EMAX Hawk 7) mass gap — non-selected frame","RETIRED","Retired 2026-08-08 from candidates.sysml (see MODEL_ISSUES.md B4). No effect on baseline."),
    ("Phase 4 (OpenHD) component pricing","TBD","Confirm before any Phase 4 build (deferred)"),
    ("Empirical validation of endurance / thermal / RF","Action","Flight test, field thermal capture, range test"),
]:
    oi.append([a,b,c])
bco = {}
for i in range(1,len(oi)):
    t = oi[i][1]
    bco[(i,1)] = AMBER if t in ("TBR","TBD") else (ACCENT if t=="Doc" else SLATE)
add_table(s, ML, Inches(1.5), CW, len(oi),
          [Inches(5.9), Inches(1.1), Inches(5.09)], oi,
          aligns=[PP_ALIGN.LEFT,PP_ALIGN.CENTER,PP_ALIGN.LEFT],
          row_h=Inches(0.62), body_size=11, header_size=11, body_colors=bco)
txt(s, ML, Inches(6.35), CW, Inches(0.5),
    [[("None of the open items block the design baseline; all are bounded and assignable. Empirical validation is the primary pre-operational activity.", 10.5, False, MUTED, True)]])

# ---- 34. CDR closure ----
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, 0, SW, Inches(0.14), AMBER)
txt(s, ML, Inches(0.55), CW, Inches(0.3), [[("SECTION 5 · CLOSURE", 12, True, RGBColor(0x8F,0xC5,0xE0))]])
txt(s, ML, Inches(0.9), CW, Inches(0.7), [[("CDR Exit Assessment & Recommendation", 30, True, WHITE)]])
# exit criteria checklist (left)
crit = [
    ("Detailed design complete & traces to R1–R8", True),
    ("Interfaces defined (internal + external)", True),
    ("Analyses show compliance with margin", True),
    ("R4 cost ✓ (26% margin) · R6 ✓ · R7 ✓ · thermal detect ✓", True),
    ("R8 stretch (60 min) not met — documented", False),
    ("Risks identified with mitigations", True),
    ("Open items bounded (none design-blocking)", True),
    ("Empirical validation pending (design gate, not yet test)", False),
]
txt(s, ML, Inches(1.95), Inches(6.6), Inches(0.35), [[("EXIT CRITERIA", 12, True, RGBColor(0x8F,0xC5,0xE0))]])
tb = s.shapes.add_textbox(ML, Inches(2.4), Inches(6.7), Inches(4.3)); tf = tb.text_frame; tf.word_wrap=True
first=True
for text, ok in crit:
    p = tf.paragraphs[0] if first else tf.add_paragraph(); first=False
    p.space_after = Pt(8); p.line_spacing=1.02
    rr = p.add_run(); rr.text = ("✓  " if ok else "▲  ")
    rr.font.name=FONT; rr.font.size=Pt(13.5); rr.font.bold=True
    rr.font.color.rgb = OKGRN if ok else AMBER
    r2 = p.add_run(); r2.text = text
    r2.font.name=FONT; r2.font.size=Pt(13.5); r2.font.color.rgb=RGBColor(0xE6,0xEC,0xF2)
# recommendation card (right)
rect(s, Inches(7.5), Inches(1.95), Inches(5.2), Inches(4.75), RGBColor(0x14,0x30,0x52), line=RGBColor(0x2a,0x44,0x63), line_w=Pt(1), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, Inches(7.8), Inches(2.2), Inches(4.6), Inches(0.4), [[("RECOMMENDATION", 12, True, AMBER)]])
txt(s, Inches(7.8), Inches(2.7), Inches(4.6), Inches(3.9),
    [[("Proceed to procurement, fabrication & integration.", 16, True, WHITE)],
     [("The detailed design is complete, fully traced, and verified by analysis with margin on cost, endurance, RF range, and thermal detection.", 12.5, False, RGBColor(0xC9,0xD8,0xE6))],
     [("Gate to operations on: (1) an instrumented Phase-1 flight to correlate the endurance model, (2) field thermal capture to validate detection/classification, and (3) an established FAA BVLOS regulatory path (Part 107).", 12.5, False, RGBColor(0xC9,0xD8,0xE6))]],
    space_after=Pt(10), line_spacing=1.06)
txt(s, ML, SH-Inches(0.5), CW, Inches(0.3),
    [[("Thermal Surveillance UAS · Critical Design Review · 2026-07-26 · Model-based (SysML v2), git-controlled", 10, False, RGBColor(0x7F,0x95,0xAB))]])

# ============================================================================
# ============================================================================
# SPEAKER NOTES  (one entry per slide, in final order)
# ============================================================================
NOTES = [
# 1 Title
"Welcome to the Critical Design Review for the thermal surveillance UAS — an autonomous, daytime wildlife-scouting drone that detects and classifies animals by their thermal signature. This review presents the detailed design, shows it traces to requirements, and demonstrates compliance through analysis. The design is model-based in SysML v2 and version-controlled in git; everything you will see is generated from that single authoritative model.",
# 2 Agenda
"The purpose of a CDR is to confirm the detailed design meets requirements and is ready to build. Scope is the committed Phase 1–3 system; the Phase 4 video downlink is shown only as a deferred capability. I'll walk requirements and traceability, architecture, interfaces, behavior, the engineering analyses, then verification, risk, and closure. The exit criteria I'm holding us to are on the final slide.",
# 3 Coverage map
"This maps our package to the standard design-review content areas so you can confirm completeness at a glance. Every required area is addressed, with the slide where you'll find it. Only verification is 'plan defined' rather than complete — that is exactly right at CDR, because test execution happens at the next gate, not this one.",
# 4 Mission / ConOps
"The mission is autonomous daytime thermal survey of wildlife over open terrain — the thermal camera detects and classifies by heat signature, which works fine in daylight, so the concept is daytime-only. The operating point sets everything downstream: 90 to 120 meters altitude, 2.23 meters per second, a 2.8 kilometer survey line in up to 4.5 meters per second of wind, with at least a 5-degree target-to-background thermal contrast. A sortie plans a route, takes off, cruises while running live onboard thermal inference, investigates and classifies detections, then returns and lands. The autonomy stance: thermal is processed onboard in real time to drive route changes — no recording and no required downlink in the committed build, with a thermal-video downlink planned as a future capability in Phase 4.",
# --- Operational Concept (OV-1) ---
"This is the OV-1 — the DoDAF High-Level Operational Concept Graphic — the one picture that tells the whole mission story. By day, the drone flies an autonomous survey pattern at 90 to 120 meters, sweeping open terrain with its thermal camera and detecting and classifying wildlife and humans by their heat signature. It processes the thermal feed onboard in real time to re-route itself, so no video downlink is required. It stays linked to a single ground operator — a laptop running QGroundControl with an ELRS control-and-telemetry link and a directional video antenna — out to a 2.8-kilometer surveillance range in up to 4.5 meters per second of wind. Every callout ties back to a driving requirement: altitude R1, speed R2, detect and classify R3, range R7. Where the block diagrams that follow show structure, this shows the concept of operations.",
# 5 Divider Requirements
"We start with requirements and traceability — the foundation every later decision has to satisfy.",
# Requirements decomposition (front of section)
"This shows how the requirements decompose by subsystem in the model. The eight mission requirements at the top break down into per-subsystem requirement packages — airframe, battery, thermal camera, computer, ground control, and autonomy — each labeled with the number of requirements it holds. Every one of those subsets a mission requirement and is satisfied by a component. This is the structural map; the full requirement text and allocations are in the three RTM slides that follow. The point is that requirements are modeled formally in SysML, not just listed in a document.",
# 6 Mission requirements
"These are the eight top-level mission requirements, taken directly from the model. R1 and R2 set the flight envelope, R3 the thermal detect-and-classify capability, R4 the twenty-five-hundred-dollar cost cap, R6 and R8 endurance, and R7 range. The side panel shows the verification method for each. At CDR, analysis is complete for all of them; demonstrations and tests are planned and gated later.",
# 7 Traceability approach
"Traceability here is bidirectional and derived from the model, not maintained by hand. Each subsystem requirement subsets its parent mission requirement, and each component carries a satisfy link to the requirements it meets. The rollup on the right shows every mission requirement decomposes and allocates to real components. The full matrix follows on the next three slides.",
# 8 RTM airframe/behavior
"This is the airframe and autonomy-behavior slice of the traceability matrix. The behavior requirements — altitude hold, speed hold, and return-to-launch on link-loss and low battery — allocate to autonomy functions and the ArduPilot firmware. Every row ties a requirement to the component responsible for meeting it.",
# 9 RTM battery/camera
"Battery and thermal-camera requirements. The battery drives the endurance, energy, and voltage requirements; the camera drives detection, resolution, field of view, and NETD. These allocate to the Battery and IRCamera parts, and the cost sub-requirements roll up to the R4 budget.",
# 10 RTM SBC/GCS
"The onboard-computer and ground-control requirements. The SBC covers compute power, thermal video input, MAVLink data exchange, and the passive-cooling constraint. The GCS covers control authority, telemetry display, range, and video display. Note that the 2.8-kilometer range requirement allocates to the antennas and RF chain — we verify that directly in the link-budget analysis.",
# (Requirements-decomposition note moved to the front of the Requirements section)
# 12 Divider Architecture
"Now the system architecture — how the design is physically decomposed and interconnected.",
# 13 Block-definition (composition) view
"This is the block-definition, or composition, view — the SysML v2 analog of a classic BDD. It shows each part definition and the typed part usages it is composed of: the Aerial Observation System is a Drone plus a GCS; the Drone is composed of the airframe, battery, anti-spark filter, thermal camera, FPV camera, GPS, onboard computer, receiver, video transmitter, UBEC, and antennas; and the GCS of the laptop, video receiver, USB capture, radio, telemetry link, ground antenna, and charger. Where the internal block diagrams that follow show how these parts connect, this view shows what the system is made of. Phase-4 OpenHD parts are omitted for the committed build.",
# 14 System context
"At the top level the system is a drone plus a laptop-based ground station, joined by three wireless links: analog FPV video on 5.8 gigahertz for piloting, ELRS on 2.4 gigahertz carrying both control and telemetry, and GNSS receive for navigation. The laptop is existing kit — modeled as an external actor and excluded from cost. Phase 4 would add a digital thermal downlink on the same 5.8 band.",
# 14 Drone IBD
"This is the detailed internal wiring of the drone, with the airframe as the hub. Follow the power chain: battery, through the anti-spark filter, into the airframe, with a separate regulated 12-volt feed to the onboard computer. The thermal camera connects to that computer over USB; the computer talks MAVLink to the flight controller to modify the route in flight. A separate analog camera feeds the video transmitter purely for piloting.",
# 15 SBC software
"Inside the onboard computer, the software decomposes into the mission app, the thermal inference model running on the RKNN neural-processing runtime, and a MAVLink router to the flight controller. USB-A port one carries the thermal feed; port two is reserved for the Phase 4 downlink; USB-C is power. The board is passive-cooled and stays under ten watts, which is how it meets the power and thermal requirements without a fan.",
# 16 GCS
"The ground station is the laptop running QGroundControl, a 5.8-gigahertz analog video receiver captured over USB, the handheld ELRS radio for manual control, and the directional patch antenna that closes the video link at range. The charger shown is bench support equipment — not flown, and not counted in system cost.",
# 17 Interfaces
"This consolidates the interface design — internal wired interfaces on the left, external wireless on the right. The important part is the bottom box: compatibility is enforced formally. The cell-count rule keeps the battery matched to the airframe, the video-format chain stays consistent end to end, RF bands match per link, and the entire power chain is standardized on XT60. These are modeled as SysML constraints and actually enforced by our configuration sweep, so an incompatible build can't slip through.",
# Standards Profile (StdV-1)
"This is the StdV-1 — the standards profile — the open and de-facto-standard protocols the system is built on. Command and telemetry are MAVLink 2; the control link is ExpressLRS over CRSF; the thermal camera is standard USB video class; navigation is u-blox UBX; the flight stack is ArduPilot; the ground station is QGroundControl; inference runs on Rockchip's RKNN. The takeaway is no proprietary lock-in on the committed build — every committed interface rides an open or de-facto standard — and the one regulatory item, FAA Part 107 plus a BVLOS authorization, is called out here too.",
# 18 Divider Behavior
"Next, the functional and behavioral design — what the system does, and how it sequences and fails safe.",
# 19 State machine
"The flight-mode state machine defines nominal operation and, more importantly, the failsafes. From the flying state, the vehicle automatically returns to launch on link-loss beyond the failsafe timeout, or when battery energy reaches the return reserve. These transitions map directly onto ArduPilot failsafe parameters. One open item to flag: the exact reserve-energy threshold value is still to be set.",
# 20 Autonomy loop
"This is the autonomous detect-investigate-classify loop that runs onboard during a sortie. Sweep-and-detect runs per-frame inference and, on a confident detection, marks the point of interest and alerts the operator. Investigate-and-classify then reroutes to the target and classifies in a loop until confident, adjusting orbit if needed, before resuming the survey. Both halves formally signal the mission context. It all runs on the onboard computer and commands routes to the flight controller over MAVLink.",
# 21 Divider Analyses
"Now the quantitative analyses that verify the design. For each one I'll state the method and, just as importantly, its fidelity — what it does and does not yet prove.",
# 22 Selected config
"These components were down-selected from about 110 real market candidates through documented trade studies. The selections are version-controlled and machine-enforced in the analysis code, so they cannot silently drift. The reference build lands at roughly 58 minutes hover and about eighteen-hundred-fifty dollars — inside every requirement. On the next two slides, here is the actual hardware.",
# 23 HW air vehicle
"Here is the actual selected flight hardware, all real purchasable parts: the iFlight Chimera9 ECO airframe, the Amprius silicon-anode Li-ion pack that drives our endurance, the PurpleRiver Mini 640 thermal core, the NanoPi M5 inference computer, the true-diversity ELRS receiver, and the anti-spark filter on the power chain. These are not concept renders — they are the components in the bill of materials.",
# 24 HW ground
"And the ground-control and support hardware: the RadioMaster TX12 handheld radio, the TrueRC directional patch antenna that closes the 2.8-kilometer link, and the HOTA dual-channel charger for the Li-ion packs. The ground kit also includes the Skydroid USB video receiver, the ELRS telemetry dongle, the airframe-bundled GPS, and the existing laptop as the ground station.",
# 25 Endurance
"Analysis one, endurance. The method is momentum, or actuator-disk, theory with a forward-flight drag term — the same physics family as eCalc. The fidelity is first-order comparative, with standard assumptions for figure-of-merit and efficiency, and — this is important — it is not yet validated against a measured hover current. The baseline predicts 58 minutes hover at only about 15 percent throttle, comfortably meeting the 30-minute requirement, though not the 60-minute stretch. Notably, the analysis is closed-loop with the model: it reads components and batteries from the model and writes ranked results back.",
# 26 Endurance vs cost
"This scatter plots every compatible configuration — total system cost against flight time, with the requirement lines overlaid. The takeaway is simple: the entire selected trade space sits to the left of the cost cap and above the 30-minute line, so cost is not the binding constraint. Our reference build sits comfortably in that sweet spot.",
# 27 Thermal
"Analysis two, thermal detection, using the Johnson criteria. We compute ground sample distance from the sensor geometry and altitude, then pixels-on-target against the detection and recognition thresholds. The optics deliver detection-grade pixels on deer and humans across the full 90-to-120-meter band. The honest caveat is on the right: robust species-level classification depends on the trained AI model, and validating that requires real field thermal capture — geometry alone does not prove it.",
# 28 RF
"Analysis three, the RF link budget. The method is free-space path loss with a 10-decibel fade-margin design target at 2.8 kilometers. All four links pass — control and telemetry with enormous margin, and the analog video link passing specifically because of the directional patch antenna, which is the enabling component. Fidelity: it is a free-space model, and the fade margin is what absorbs foliage and multipath. It is not yet range-tested.",
# 29 Physical fit
"Analysis four asks the practical question: does it physically fit? We model each airframe's usable deck and each component's envelope. The binding constraint is the NanoPi M5's 90-by-62-millimeter footprint. The selected Chimera9 deck accepts it with about 8 millimeters to spare — and that fit is a key reason we chose it over the lighter endurance leader, which is a no-fit. Caveat: the deck dimensions are still estimates and should be confirmed against vendor CAD before the Phase 2 build.",
# 30 Cost
"Analysis five, cost. The committed system rolls up to about eighteen-hundred-forty-eight dollars against the twenty-five-hundred-dollar cap — roughly 26 percent margin. The thermal camera at six-hundred-fifty dollars is the single largest line. Even all four phases together stay around two thousand dollars. Cost is a solved problem for this design.",
# Performance measures (SV-7)
"This is the SV-7 — the systems-measures scorecard — the quantitative bottom line of the analysis section. Each row pairs a performance measure with its requirement threshold, the as-designed value from the analysis, the margin, and a pass or fail. The headline is that every hard requirement is met with margin: endurance nearly double the thirty-minute floor, range and RF margin positive, thermal pixels-on-target above the Johnson thresholds, and cost twenty-six percent under the cap. The only miss is the sixty-minute stretch goal, which we flag rather than hide, and SBC power sits right at its ten-watt limit.",
# 31 Divider Closure
"Finally: verification planning, the manufacturing approach, risk, open items, and the CDR exit assessment.",
# 32 VCRM
"The verification cross-reference matrix shows how every requirement will be verified and its status today. Cost and soldering are verified now, by inspection. Everything else is analysis-complete with test or demonstration planned — which is the correct posture at CDR; test readiness is the next gate. The one exception is the 60-minute stretch goal, which the design does not meet, and we have documented that rather than hidden it.",
# 33 Manufacturing
"Producibility. We minimize soldering, per requirement R5, by using a plug-and-play airframe — the only solder joint in the committed build is the onboard computer's power leads. Nearly everything is off-the-shelf; the single fabricated part is a 3D-printed computer mount. And the build is phased into three incremental stages, each independently testable before the next is added.",
# Project timeline (PV-2)
"This is the PV-2 — the project timeline — the roadmap in one picture. Three committed phases build on each other: Phase 1 is basic flight with FPV and waypoints; Phase 2 adds the thermal payload and onboard computer; Phase 3 is software-only, the AI detection and autonomous re-routing. Phase 4, in grey, is the deferred digital video downlink. The banner marks where we are — at CDR, detailed design complete, with procurement as the next gate. There are no calendar dates yet; this shows sequence and dependency, and that each phase is independently testable before the next.",
# 34 Risk
"The risk posture. The three highest risks are: first, the design is modeled but not yet flown, so the headline numbers are predictions; second, species classification is unproven without real thermal data; and third, the regulatory path — at 2.8 kilometers the survey is beyond visual line of sight, which requires FAA authorization, and commercial operation requires a Part 107 certificate. Operating in daylight rather than after dark lowers this risk, but the BVLOS range keeps it on the list. Each has a defined mitigation. At CDR, 'not yet built' is expected; the material risks are the estimates and these external items, and we are naming them openly.",
# 35 Open items
"Our open items are bounded and assignable — none of them block the design baseline. They include setting the return-to-launch reserve threshold, reconciling one stale entry in the issues log, confirming the selected frame's real deck dimensions, and running the empirical validation campaign. That validation — flight test, field thermal capture, and range test — is the primary pre-operational activity.",
# 36 Closure
"In summary: the detailed design is complete, fully traced, and verified by analysis with margin on cost, endurance, RF range, and thermal detection. The recommendation is to proceed to fabrication and integration, gated to operations on three things — an instrumented flight to validate the endurance model, field thermal capture to validate detection and classification, and an established regulatory path. Thank you; I'll take your questions.",
]
assert len(NOTES) == len(prs.slides._sldIdLst), ("notes/slide mismatch", len(NOTES), len(prs.slides._sldIdLst))
for sl, note in zip(prs.slides, NOTES):
    sl.notes_slide.notes_text_frame.text = note

prs.save(OUT)
print("SAVED:", OUT)
print("Total slides:", len(prs.slides._sldIdLst))
print("Notes attached:", len(NOTES))

